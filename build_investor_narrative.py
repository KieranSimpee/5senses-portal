from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette (exact from template) ──────────────────────────────────────
VIOLET      = RGBColor(0x5E, 0x50, 0xFB)   # #5E50FB  accent
LILAC       = RGBColor(0x8C, 0x82, 0xFC)   # #8C82FC  secondary
LAVENDER    = RGBColor(0xE8, 0xE6, 0xFE)   # #E8E6FE  background wash
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1F)   # #1A1A1F  body text
GREY        = RGBColor(0x5A, 0x5A, 0x6A)   # #5A5A6A  sub text
LIGHTGREY   = RGBColor(0x99, 0x99, 0xAA)   # #9999AA  labels

FONT = "Montserrat"

W  = Inches(13.33)
H  = Inches(7.5)
LB = Inches(0.35)   # left bar width
RB = Inches(0.18)   # right bar width

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helper: blank slide ────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

# ── Helper: add rectangle ──────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ── Helper: add textbox ────────────────────────────────────────────────────────
def add_text(slide, text, l, t, w, h,
             font_name=FONT, font_size=12, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return txb

# ── Helper: left violet bar + right bar (standard inner slides) ────────────────
def std_bars(slide):
    add_rect(slide, 0,        0, LB,        H, VIOLET)
    add_rect(slide, W - RB,   0, RB,        H, LILAC)

# ── Helper: section label (top left after bar) ────────────────────────────────
def section_label(slide, text):
    add_text(slide, text,
             LB + Inches(0.25), Inches(0.28), Inches(5), Inches(0.3),
             font_size=9, bold=True, color=LIGHTGREY)

# ── Helper: headline ──────────────────────────────────────────────────────────
def headline(slide, text, font_size=32, color=BLACK, top=Inches(0.7)):
    add_text(slide, text,
             LB + Inches(0.25), top, Inches(11.5), Inches(1.2),
             font_size=font_size, bold=True, color=color)

# ── Helper: sub-headline ──────────────────────────────────────────────────────
def sub_headline(slide, text, top=Inches(1.6)):
    add_text(slide, text,
             LB + Inches(0.25), top, Inches(10), Inches(0.5),
             font_size=15, bold=False, color=GREY)

# ── Helper: bullet card (lavender bg + violet left stripe) ────────────────────
def bullet_card(slide, title, body, left, top, width=Inches(2.8), height=Inches(1.5)):
    add_rect(slide, left, top, width, height, LAVENDER)
    add_rect(slide, left, top, Inches(0.07), height, VIOLET)
    add_text(slide, title,
             left + Inches(0.18), top + Inches(0.15), width - Inches(0.25), Inches(0.4),
             font_size=13, bold=True, color=VIOLET)
    add_text(slide, body,
             left + Inches(0.18), top + Inches(0.55), width - Inches(0.25), height - Inches(0.65),
             font_size=11, bold=False, color=BLACK)

# ── Helper: step card (white bg + violet left stripe) ─────────────────────────
def step_card(slide, num, title, body, left, top, width=Inches(2.8), height=Inches(2.4)):
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.07), height, VIOLET)
    add_text(slide, num,
             left + Inches(0.18), top + Inches(0.1), Inches(1), Inches(0.7),
             font_size=26, bold=True, color=LILAC)
    add_text(slide, title,
             left + Inches(0.18), top + Inches(0.75), width - Inches(0.25), Inches(0.45),
             font_size=13, bold=True, color=BLACK)
    add_text(slide, body,
             left + Inches(0.18), top + Inches(1.2), width - Inches(0.25), height - Inches(1.3),
             font_size=11, bold=False, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)

# Background blocks (matching template cover layout)
add_rect(s1, 0, 0, LB, H, VIOLET)                           # left violet bar
add_rect(s1, LB, 0, Inches(8.5), H, WHITE)                  # white main area
add_rect(s1, Inches(8.85), 0, Inches(4.48), H, LAVENDER)    # right lavender panel
add_rect(s1, W - RB, 0, RB, H, LILAC)                        # right lilac bar

# Bottom violet strip
add_rect(s1, 0, H - Inches(0.55), W, Inches(0.55), VIOLET)

# Bottom strip text
add_text(s1, "ONE STOP ASIAN BEAUTY  :  SIMPLIFY TO AMPLIFY",
         LB + Inches(0.3), H - Inches(0.45), Inches(9), Inches(0.38),
         font_size=10, bold=True, color=WHITE)
add_text(s1, "kieran@simplex-ity.com",
         Inches(10.5), H - Inches(0.45), Inches(2.5), Inches(0.38),
         font_size=10, bold=False, color=LIGHTGREY)

# Deck label top
add_text(s1, "INVESTOR NARRATIVE DECK  ·  2026",
         LB + Inches(0.35), Inches(0.3), Inches(7), Inches(0.35),
         font_size=10, bold=True, color=LIGHTGREY)

# Main headline
add_text(s1, "SIMPLEX-ITY",
         LB + Inches(0.35), Inches(1.2), Inches(8), Inches(1.2),
         font_size=52, bold=True, color=BLACK)
add_text(s1, "Simplify to Amplify.",
         LB + Inches(0.35), Inches(2.3), Inches(8), Inches(0.9),
         font_size=42, bold=True, color=VIOLET)
add_text(s1, "Investor Narrative 2026",
         LB + Inches(0.35), Inches(3.25), Inches(7), Inches(0.5),
         font_size=18, bold=False, color=GREY)

# Right panel accent text
add_text(s1, "AI in Reality.\nROI in Action.",
         Inches(9.1), Inches(2.2), Inches(3.8), Inches(1.8),
         font_size=22, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
add_text(s1, "Partner-backed.\nData-driven.\nInvestable.",
         Inches(9.1), Inches(4.2), Inches(3.8), Inches(1.5),
         font_size=13, bold=False, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HONEST REFLECTION
# ═══════════════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
std_bars(s2)
section_label(s2, "01 — HONEST REFLECTION")
headline(s2, "Where we started — and what we learned.", font_size=30)
sub_headline(s2, "A candid self-assessment. Transparency builds trust.", top=Inches(1.55))

bullet_card(s2, "No Detailed Plan",
            "Founded without a structured expansion roadmap or go-to-market strategy.",
            LB + Inches(0.25), Inches(2.1), Inches(2.9), Inches(1.55))
bullet_card(s2, "Underestimated Difficulty",
            "Assumed platform-building was simple. The true ambition was always high-impact — fundamentally harder.",
            LB + Inches(3.35), Inches(2.1), Inches(2.9), Inches(1.55))
bullet_card(s2, "Vision vs Reality Gap",
            "Resources, support, and capacity were never reconciled with the scale of the vision.",
            LB + Inches(6.45), Inches(2.1), Inches(2.9), Inches(1.55))
bullet_card(s2, "Over-Optimism as a Blind Spot",
            "Hard work alone was not enough. The internal vision was never made visible to others.",
            LB + Inches(9.55), Inches(2.1), Inches(2.9), Inches(1.55))

# Mantra box
add_rect(s2, LB + Inches(0.25), Inches(4.1), Inches(12.2), Inches(0.9), VIOLET)
add_text(s2, "Mantra:  \"Build first. Prove second. Invite third.\"",
         LB + Inches(0.55), Inches(4.2), Inches(11.5), Inches(0.7),
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SHIFT
# ═══════════════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
std_bars(s3)
section_label(s3, "02 — THE SHIFT")
headline(s3, "Foundation before invitation.", font_size=34)
sub_headline(s3, "Leadership is not convincing people to believe. It is building something they cannot doubt.", top=Inches(1.58))

# Three flow cards with arrows
card_w = Inches(3.3)
card_h = Inches(2.2)
tops   = Inches(2.2)

# Card 1 – Foundation
step_card(s3, "01", "Foundation",
          "Build AI partners, government backing, and a live demonstrable platform before asking anyone to join.",
          LB + Inches(0.25), tops, card_w, card_h)

# Arrow 1
add_text(s3, "→", Inches(4.2), tops + Inches(0.7), Inches(0.5), Inches(0.7),
         font_size=22, bold=False, color=LILAC)

# Card 2 – Proof
step_card(s3, "02", "Proof",
          "3-month trial with real brands and influencers. Live data, consumer behaviour, real results.",
          LB + Inches(4.1), tops, card_w, card_h)

# Arrow 2
add_text(s3, "→", Inches(7.95), tops + Inches(0.7), Inches(0.5), Inches(0.7),
         font_size=22, bold=False, color=LILAC)

# Card 3 – Invitation
step_card(s3, "03", "Invitation",
          "With data, investment, and a live platform — invite trusted people to join on solid ground.",
          LB + Inches(7.95), tops, card_w, card_h)

# Key insight
add_rect(s3, LB + Inches(0.25), Inches(4.85), Inches(12.2), Inches(0.75), LAVENDER)
add_text(s3, "Key Insight:  Others' hesitation was never a lack of courage — it was a rational response to an insufficient foundation.",
         LB + Inches(0.45), Inches(4.95), Inches(11.8), Inches(0.6),
         font_size=13, bold=False, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TANGIBLE PROOF SO FAR
# ═══════════════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
std_bars(s4)
section_label(s4, "03 — TANGIBLE PROOF SO FAR")
headline(s4, "This is no longer just a vision.", font_size=34)
sub_headline(s4, "Proof of concept is live. AI ecosystem is being built.", top=Inches(1.58))

bullet_card(s4, "Tint AI Try-On — LIVE",
            "Virtual makeup try-on integrated and demonstrable directly on the platform. Real AI, real product.",
            LB + Inches(0.25), Inches(2.1), Inches(3.9), Inches(1.7))
bullet_card(s4, "Korean Skincare Connections",
            "Early brand relationships established through direct sourcing background. Pipeline is warm.",
            LB + Inches(4.4), Inches(2.1), Inches(3.9), Inches(1.7))
bullet_card(s4, "AI Partner Ecosystem",
            "Tint AI confirmed. Vybd AI Commerce being onboarded. Partners are real, not aspirational.",
            LB + Inches(8.55), Inches(2.1), Inches(3.9), Inches(1.7))

# Stats row
add_rect(s4, LB + Inches(0.25), Inches(4.2), Inches(3.9), Inches(1.55), LAVENDER)
add_rect(s4, LB + Inches(0.25), Inches(4.2), Inches(0.07), Inches(1.55), VIOLET)
add_text(s4, "Proof of Concept", LB + Inches(0.45), Inches(4.3), Inches(3.5), Inches(0.4), font_size=11, bold=True, color=VIOLET)
add_text(s4, "Tint AI integrated & live", LB + Inches(0.45), Inches(4.72), Inches(3.5), Inches(0.35), font_size=13, bold=True, color=BLACK)

add_rect(s4, LB + Inches(4.4), Inches(4.2), Inches(3.9), Inches(1.55), LAVENDER)
add_rect(s4, LB + Inches(4.4), Inches(4.2), Inches(0.07), Inches(1.55), VIOLET)
add_text(s4, "Market Focus", LB + Inches(4.6), Inches(4.3), Inches(3.5), Inches(0.4), font_size=11, bold=True, color=VIOLET)
add_text(s4, "US & Canada — Korean beauty", LB + Inches(4.6), Inches(4.72), Inches(3.5), Inches(0.35), font_size=13, bold=True, color=BLACK)

add_rect(s4, LB + Inches(8.55), Inches(4.2), Inches(3.9), Inches(1.55), LAVENDER)
add_rect(s4, LB + Inches(8.55), Inches(4.2), Inches(0.07), Inches(1.55), VIOLET)
add_text(s4, "Platform Status", LB + Inches(8.75), Inches(4.3), Inches(3.5), Inches(0.4), font_size=11, bold=True, color=VIOLET)
add_text(s4, "Pre-launch · Trial ready", LB + Inches(8.75), Inches(4.72), Inches(3.5), Inches(0.35), font_size=13, bold=True, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROADMAP FORWARD
# ═══════════════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
std_bars(s5)
section_label(s5, "04 — ROADMAP FORWARD")
headline(s5, "The correct order — executed with discipline.", font_size=32)
sub_headline(s5, "Five sequential steps from where we are to investment-ready.", top=Inches(1.58))

steps = [
    ("01", "Government Support", "StartmeUp HK engagement for grants, credibility, and official backing."),
    ("02", "Operational Partner", "Vybd AI Commerce onboard — US operations, logistics, FDA compliance offloaded."),
    ("03", "3-Month Trial",       "Zero-cost brand & influencer trial. 30% discount. Real data. Real feedback."),
    ("04", "NEST VC Pitch",       "Present with live platform, real data, and proven AI ecosystem to NEST VC."),
]
x_start = LB + Inches(0.25)
x_gap   = Inches(3.05)
for idx, (num, title, body) in enumerate(steps):
    step_card(s5, num, title, body, x_start + idx * x_gap, Inches(2.05), Inches(2.9), Inches(2.6))
    if idx < 3:
        add_text(s5, "→", x_start + idx * x_gap + Inches(2.95), Inches(2.75),
                 Inches(0.3), Inches(0.6), font_size=20, bold=False, color=LILAC)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRIAL LAUNCH MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
std_bars(s6)
section_label(s6, "05 — TRIAL LAUNCH MODEL")
headline(s6, "3 months. Zero cost. Real results.", font_size=34)
sub_headline(s6, "Authentic market positioning — this is a trial, and we say so.", top=Inches(1.58))

trial_points = [
    ("Zero-Cost Entry",      "Brands and influencers participate at no charge during the 3-month trial window."),
    ("30% Consumer Discount","All platform products priced at 30% below market to drive trial traffic and purchases."),
    ("Honest Positioning",   "Consumers are told this is a trial. Authentic, transparent — builds real trust."),
    ("Platform Debugging",   "Use 3 months to stress-test systems, fix issues, and improve the full user experience."),
]
for i, (title, body) in enumerate(trial_points):
    col = i % 2
    row = i // 2
    lft = LB + Inches(0.25) + col * Inches(6.2)
    top = Inches(2.1) + row * Inches(1.75)
    bullet_card(s6, title, body, lft, top, Inches(5.95), Inches(1.55))

# Revenue row
add_rect(s6, LB + Inches(0.25), Inches(5.75), Inches(5.9), Inches(1.0), VIOLET)
add_text(s6, "Revenue: Transaction commissions + consumer behaviour intelligence data",
         LB + Inches(0.45), Inches(5.88), Inches(5.5), Inches(0.75),
         font_size=13, bold=True, color=WHITE)

add_rect(s6, LB + Inches(6.45), Inches(5.75), Inches(6.05), Inches(1.0), LILAC)
add_text(s6, "Long-term: ABW/YesStyle becomes a vendor on our terms — not a dependency",
         LB + Inches(6.65), Inches(5.88), Inches(5.65), Inches(0.75),
         font_size=13, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CAPITAL NEED
# ═══════════════════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
std_bars(s7)
section_label(s7, "06 — CAPITAL NEED")
headline(s7, "What we need — and why it pays back.", font_size=34)
sub_headline(s7, "One clear funding purpose. One clear return path.", top=Inches(1.58))

# Main capital box
add_rect(s7, LB + Inches(0.25), Inches(2.1), Inches(12.2), Inches(1.5), LAVENDER)
add_rect(s7, LB + Inches(0.25), Inches(2.1), Inches(0.1), Inches(1.5), VIOLET)
add_text(s7, "Primary Use of Funds",
         LB + Inches(0.55), Inches(2.2), Inches(5), Inches(0.4),
         font_size=11, bold=True, color=VIOLET)
add_text(s7, "Cover Tint AI Try-On click costs during the 3-month trial period",
         LB + Inches(0.55), Inches(2.62), Inches(11.5), Inches(0.75),
         font_size=20, bold=True, color=BLACK)

# ROI path cards
bullet_card(s7, "Real Consumer Data",
            "Every AI try-on click generates purchase intent signals, preferences, and conversion data. This is our dataset for NEST VC.",
            LB + Inches(0.25), Inches(3.95), Inches(3.9), Inches(1.85))
bullet_card(s7, "Platform Validation",
            "Trial data proves the model works — conversion rates, influencer reliability scores, and brand satisfaction documented.",
            LB + Inches(4.4), Inches(3.95), Inches(3.9), Inches(1.85))
bullet_card(s7, "Investor Confidence",
            "NEST VC receives a pitch backed by real numbers from a live platform — not a pitch deck promise.",
            LB + Inches(8.55), Inches(3.95), Inches(3.9), Inches(1.85))

add_rect(s7, LB + Inches(0.25), Inches(6.1), Inches(12.2), Inches(0.65), VIOLET)
add_text(s7, "Clear ROI Path: Trial data → NEST VC validation → Investment secured → Platform scales",
         LB + Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THE ASK
# ═══════════════════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)

# Full lavender background
add_rect(s8, 0, 0, W, H, LAVENDER)
# Left violet bar
add_rect(s8, 0, 0, LB, H, VIOLET)
# Right bar
add_rect(s8, W - RB, 0, RB, H, LILAC)
# Bottom strip
add_rect(s8, 0, H - Inches(0.55), W, Inches(0.55), VIOLET)
add_text(s8, "ONE STOP ASIAN BEAUTY  :  SIMPLIFY TO AMPLIFY",
         LB + Inches(0.3), H - Inches(0.45), Inches(9), Inches(0.38),
         font_size=10, bold=True, color=WHITE)
add_text(s8, "kieran@simplex-ity.com",
         Inches(10.5), H - Inches(0.45), Inches(2.5), Inches(0.38),
         font_size=10, bold=False, color=LIGHTGREY)

section_label(s8, "07 — THE ASK")

# Central headline
add_text(s8, "Tangible. Investable. Partner-ready.",
         LB + Inches(0.35), Inches(1.0), Inches(12), Inches(1.1),
         font_size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Three pillars
pillar_w = Inches(3.7)
pillar_h = Inches(1.8)
pillar_t = Inches(2.35)
for i, (title, body) in enumerate([
    ("Reduced Risk",      "AI ecosystem live. Trial data collected. Nothing is theoretical anymore."),
    ("Increased Credibility", "Government support + AI partners + real brand connections = fundable foundation."),
    ("Proven Model",      "3-month trial proves platform viability before we ask for serious capital."),
]):
    lft = LB + Inches(0.25) + i * Inches(4.05)
    add_rect(s8, lft, pillar_t, pillar_w, pillar_h, WHITE)
    add_rect(s8, lft, pillar_t, Inches(0.07), pillar_h, VIOLET)
    add_text(s8, title, lft + Inches(0.18), pillar_t + Inches(0.15), pillar_w - Inches(0.25), Inches(0.45),
             font_size=14, bold=True, color=VIOLET)
    add_text(s8, body,  lft + Inches(0.18), pillar_t + Inches(0.62), pillar_w - Inches(0.25), pillar_h - Inches(0.75),
             font_size=12, bold=False, color=BLACK)

# Closing statement
add_rect(s8, LB + Inches(0.25), Inches(4.5), Inches(12.2), Inches(1.35), VIOLET)
add_text(s8, "\"AI in Reality, ROI in Action.\"",
         LB + Inches(0.35), Inches(4.65), Inches(12), Inches(0.65),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s8, "SIMPLEX-ITY  ·  Investor Narrative 2026  ·  kieran@simplex-ity.com  ·  www.simplex-ity.com",
         LB + Inches(0.35), Inches(5.2), Inches(12), Inches(0.45),
         font_size=11, bold=False, color=LAVENDER, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
OUT = "/app/SIMPLEX-ITY_Investor_Narrative_2026.pptx"
prs.save(OUT)
print("Saved:", OUT)

import os
print("Size:", os.path.getsize(OUT), "bytes")
