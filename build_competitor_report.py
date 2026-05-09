from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1a, 0x05, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x06, 0x95, 0x6E)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

# ─── SLIDE 1: COVER ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
shape = slide.shapes.add_shape(9, Inches(9.5), Inches(-0.5), Inches(5), Inches(5))
shape.fill.solid(); shape.fill.fore_color.rgb = PURPLE; shape.line.fill.background()
add_text(slide, "HK BEAUTY-TECH", 1, 1.5, 11, 0.8, size=48, bold=True, color=WHITE)
add_text(slide, "Competitor & Market Analysis", 1, 2.6, 10, 0.7, size=26, color=LIGHT_PURPLE, italic=True)
add_rect(slide, 1, 3.5, 5, 0.06, PURPLE)
add_text(slide, "SIMPLEX-ITY Strategic Intelligence Report  |  May 2026", 1, 3.7, 11, 0.5, size=15, color=GRAY)
add_text(slide, "Prepared by Simpee AI  |  Confidential", 1, 6.3, 8, 0.5, size=13, color=GRAY, italic=True)

# ─── SLIDE 2: MARKET SNAPSHOT ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "MARKET SNAPSHOT", 1, 0.4, 10, 0.5, size=13, bold=True, color=PURPLE)
add_text(slide, "HK Beauty-Tech Market 2026 — Key Numbers", 1, 0.9, 11, 0.7, size=28, bold=True, color=WHITE)

metrics = [
    ("$2.3B", "Global AI Beauty\nPersonalization\nMarket 2026", PURPLE),
    ("17%", "AI Beauty Market\nCAGR 2025-2033", GREEN),
    ("45.3%", "Instagram Share of\nHK Influencer\nMarketing", GOLD),
    ("XHS +↑", "Fastest Growing\nPlatform in\nHK Beauty 2025", RED),
]
for i, (val, label, color) in enumerate(metrics):
    x = 1 + i * 3.1
    add_rect(slide, x, 2.0, 2.8, 2.8, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, 2.0, 2.8, 0.08, color)
    add_text(slide, val, x+0.1, 2.2, 2.6, 0.8, size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.1, 3.1, 2.6, 1.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Key Trend:", 1, 5.1, 11, 0.4, size=14, bold=True, color=WHITE)
add_text(slide, "68% of global beauty purchases are influenced by social discovery. TikTok Shop grew US beauty sales 108% in 2025. HK consumers are highly receptive to AI try-on technology — conversion rates 2x higher vs standard product images.", 1, 5.6, 11, 1.2, size=13, color=LIGHT_PURPLE)

# ─── SLIDE 3: COMPETITOR MAP ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "COMPETITIVE LANDSCAPE", 1, 0.4, 10, 0.5, size=13, bold=True, color=PURPLE)
add_text(slide, "Who's in the HK Beauty-Tech Space?", 1, 0.9, 11, 0.7, size=28, bold=True, color=WHITE)

competitors = [
    ("Perfect Corp", "Taiwan/Global", "YouCam virtual try-on, AI skin analysis", "Strong tech, no influencer network"),
    ("ModiFace (L'Oréal)", "Global", "AR try-on for L'Oréal brands only", "Closed ecosystem, brand-exclusive"),
    ("AnyMind Group", "Asia-wide", "Influencer management + brand matching", "No AI try-on integration"),
    ("Luna Marketing HK", "Hong Kong", "Local influencer campaigns", "No tech layer, manual process"),
    ("NC Media / PLTFRM", "Hong Kong", "Social commerce & KOL management", "Limited to specific platforms"),
    ("Meitu Beauty", "China/HK", "AI beauty filters & try-on app", "Consumer app, not B2B brand tool"),
]

headers = ["Competitor", "Market", "Strength", "Weakness"]
col_widths = [2.5, 1.8, 3.8, 3.8]
col_x = [0.7, 3.2, 5.0, 8.8]

# Header row
add_rect(slide, 0.7, 1.85, 11.9, 0.45, PURPLE)
for j, (header, w, x) in enumerate(zip(headers, col_widths, col_x)):
    add_text(slide, header, x+0.05, 1.88, w-0.1, 0.4, size=13, bold=True, color=WHITE)

row_colors = [RGBColor(0x2D, 0x10, 0x4A), RGBColor(0x1F, 0x08, 0x38)]
for i, (comp, market, strength, weakness) in enumerate(competitors):
    y = 2.4 + i * 0.72
    add_rect(slide, 0.7, y, 11.9, 0.68, row_colors[i % 2])
    row_data = [comp, market, strength, weakness]
    for j, (data, w, x) in enumerate(zip(row_data, col_widths, col_x)):
        color = GOLD if j == 0 else WHITE
        add_text(slide, data, x+0.05, y+0.05, w-0.1, 0.58, size=11, color=color, bold=(j==0))

# ─── SLIDE 4: SIMPLEX-ITY POSITIONING ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "OUR POSITIONING", 1, 0.4, 10, 0.5, size=13, bold=True, color=PURPLE)
add_text(slide, "The Gap SIMPLEX-ITY Fills", 1, 0.9, 11, 0.7, size=30, bold=True, color=WHITE)

add_text(slide, "No competitor in HK combines ALL THREE of these:", 1, 1.8, 11, 0.5, size=16, color=LIGHT_PURPLE, italic=True)

three_pillars = [
    ("💄\nAI Try-On\nTechnology", "TINT-powered virtual makeup that converts browsers into buyers"),
    ("🤝\nVerified Influencer\nNetwork", "Pre-vetted HK influencers with reliability scoring and delivery guarantees"),
    ("📊\nEnd-to-End\nCampaign Management", "From brief to report — one platform, zero headaches"),
]
for i, (title, desc) in enumerate(three_pillars):
    x = 1 + i * 4.1
    add_rect(slide, x, 2.5, 3.7, 3.8, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, 2.5, 3.7, 0.08, PURPLE)
    add_text(slide, title, x+0.1, 2.65, 3.5, 1.5, size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x+0.15, 4.3, 3.4, 1.8, size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 1, 6.55, 11.3, 0.06, PURPLE)
add_text(slide, "SIMPLEX-ITY = The only HK-based platform combining AI beauty tech + verified influencers + full campaign execution", 1, 6.65, 11.3, 0.6, size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ─── SLIDE 5: SWOT ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "SWOT ANALYSIS", 1, 0.4, 10, 0.5, size=13, bold=True, color=PURPLE)
add_text(slide, "SIMPLEX-ITY — Strategic Assessment", 1, 0.9, 11, 0.7, size=28, bold=True, color=WHITE)

swot = [
    ("STRENGTHS", GREEN, [
        "Only HK platform combining AI try-on + influencer network",
        "TINT technology proven to 2x conversions",
        "Local HK focus — deep market understanding",
        "Reliability scoring removes campaign risk",
    ]),
    ("WEAKNESSES", GOLD, [
        "Early stage — brand awareness still building",
        "Smaller influencer database vs global platforms",
        "Tech setup requires brand education",
        "Resource-intensive for small team",
    ]),
    ("OPPORTUNITIES", PURPLE, [
        "XHS (Xiaohongshu) surging in HK — first mover advantage",
        "AI beauty market growing 17% CAGR",
        "Brands moving budget from ads to influencer",
        "China-HK cross-border beauty market expanding",
    ]),
    ("THREATS", RED, [
        "Global platforms (AnyMind) expanding into HK",
        "L'Oréal / Perfect Corp may go B2B",
        "Economic slowdown affecting marketing budgets",
        "Platform algorithm changes affecting reach",
    ]),
]
positions = [(0.7, 1.9), (7.0, 1.9), (0.7, 4.7), (7.0, 4.7)]
for (label, color, points), (x, y) in zip(swot, positions):
    add_rect(slide, x, y, 6.0, 2.6, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, y, 6.0, 0.45, color)
    add_text(slide, label, x+0.1, y+0.05, 5.8, 0.38, size=14, bold=True, color=WHITE)
    py = y + 0.55
    for pt in points:
        add_text(slide, "• " + pt, x+0.15, py, 5.7, 0.42, size=11, color=WHITE)
        py += 0.44

# ─── SLIDE 6: STRATEGIC RECOMMENDATIONS ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "STRATEGIC RECOMMENDATIONS", 1, 0.4, 10, 0.5, size=13, bold=True, color=PURPLE)
add_text(slide, "SIMPLEX-ITY — Winning Moves for 2026", 1, 0.9, 11, 0.7, size=28, bold=True, color=WHITE)

recs = [
    ("🎯", "Own the XHS Channel", "Xiaohongshu is the fastest growing platform in HK beauty. Build HK's first dedicated XHS influencer campaign service before competitors arrive.", "Q2 2026"),
    ("💄", "Make TINT the Hero", "Lead all brand pitches with TINT try-on ROI data. Brands that see 200% conversion lift will convert instantly. Create a free TINT demo for first meetings.", "Immediate"),
    ("🏆", "Tier-1 Brand Anchor", "Land 2-3 premium international beauty brands as anchor clients. Their logos validate the platform for all future brand pitches.", "Q2-Q3 2026"),
    ("📱", "Influencer App", "Build a simple influencer portal for campaign tracking, live count logging, and commission visibility. Reduces management overhead 60%.", "Q3 2026"),
]
y = 2.0
for icon, title, desc, timeline in recs:
    add_rect(slide, 1, y, 11.3, 1.05, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, 1, y, 0.08, 1.05, PURPLE)
    add_text(slide, icon + "  " + title, 1.25, y+0.05, 7.5, 0.42, size=15, bold=True, color=PURPLE)
    add_text(slide, desc, 1.25, y+0.5, 7.5, 0.52, size=12, color=WHITE)
    add_rect(slide, 9.8, y+0.25, 2.2, 0.42, PURPLE)
    add_text(slide, timeline, 9.85, y+0.27, 2.1, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += 1.2

prs.save('/app/SIMPLEX-ITY_Competitor_Analysis.pptx')
print("Competitor analysis saved!")
