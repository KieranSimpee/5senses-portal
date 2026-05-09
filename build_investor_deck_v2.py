"""
SIMPLEX-ITY Investor Pitch Deck v2
Brand-compliant: Exo 2 (headlines) + Montserrat (body)
Official palette from Design System v1.0
Logo: full-colour on white slides, text-only label on dark slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ── Brand Colors ──────────────────────────────────────────────────
LILAC      = RGBColor(0x8c, 0x82, 0xfc)
VIOLET     = RGBColor(0x5e, 0x50, 0xfb)
SOFT_LILAC = RGBColor(0xba, 0xb4, 0xfd)
LAV_WASH   = RGBColor(0xe8, 0xe6, 0xfe)
WHITE      = RGBColor(0xff, 0xff, 0xff)
NEUTRAL    = RGBColor(0xe6, 0xe6, 0xe6)
BODY       = RGBColor(0x1a, 0x1a, 0x1f)
SECONDARY  = RGBColor(0x5a, 0x5a, 0x6a)
MUTED      = RGBColor(0x99, 0x99, 0xaa)
DARK_HERO  = RGBColor(0x14, 0x10, 0x2a)
DARK2      = RGBColor(0x20, 0x14, 0x40)

# ── Fonts (per Design System v1.0) ───────────────────────────────
EXO2        = 'Exo 2'        # Primary: headlines, brand moments
MONTSERRAT  = 'Montserrat'   # Secondary: body, UI, labels

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

LOGO = '/app/logo.jpg'

# ── Helpers ───────────────────────────────────────────────────────
def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False,
        font=MONTSERRAT, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size    = Pt(size)
    r.font.bold    = bold
    r.font.color.rgb = color
    r.font.italic  = italic
    r.font.name    = font
    return tb

def logo(slide, l, t, w=2.4):
    """Full-colour logo — use only on white/very light BGs per brand rules."""
    slide.shapes.add_picture(LOGO, Inches(l), Inches(t), width=Inches(w))

def logo_text(slide, l, t, size=13, color=WHITE):
    """White text-only wordmark for dark backgrounds."""
    txt(slide, 'SIMPLEX-ITY', l, t, 3, 0.45,
        size=size, bold=True, color=color, font=EXO2)

def bar(slide, l, t, w, h=0.055, color=VIOLET):
    rect(slide, l, t, w, h, color)

def eyebrow(slide, label, l=0.35, t=0.2):
    txt(slide, label, l, t, 8, 0.28,
        size=9, bold=True, color=MUTED, font=MONTSERRAT)

def headline(slide, text, l=0.35, t=0.62, w=12.6, size=34,
             color=BODY, font=EXO2):
    txt(slide, text, l, t, w, 0.75,
        size=size, bold=True, color=color, font=font)

def sub(slide, text, l=0.35, t=1.46, w=11, size=15,
        color=SECONDARY):
    txt(slide, text, l, t, w, 0.42,
        size=size, color=color, italic=True, font=MONTSERRAT)

def stat_card(slide, val, label, x, y, w=3.05, h=2.1,
              val_color=VIOLET, bg_col=LAV_WASH, accent=VIOLET):
    rect(slide, x, y, w, h, bg_col)
    bar(slide, x, y, w, color=accent)
    txt(slide, val, x+0.1, y+0.15, w-0.2, 0.85,
        size=32, bold=True, color=val_color,
        align=PP_ALIGN.CENTER, font=EXO2)
    txt(slide, label, x+0.1, y+1.05, w-0.2, 0.9,
        size=12, color=BODY, align=PP_ALIGN.CENTER,
        font=MONTSERRAT)

def feature_card(slide, title, body, x, y, w=6.2, h=2.15,
                 accent=VIOLET, title_color=None):
    if title_color is None:
        title_color = accent
    rect(slide, x, y, w, h, LAV_WASH)
    bar(slide, x, y, w, color=accent)
    txt(slide, title, x+0.2, y+0.12, w-0.35, 0.45,
        size=14, bold=True, color=title_color, font=EXO2)
    txt(slide, body, x+0.2, y+0.63, w-0.35, h-0.75,
        size=12.5, color=BODY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER  (dark hero, white wordmark)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK_HERO)
bar(sl, 0, 0, 13.33, color=VIOLET)
rect(sl, 7.8, 0, 5.53, 7.5, DARK2)
bar(sl, 7.8, 0, 0.06, h=7.5, color=LILAC)

# White logo on dark bg — use text wordmark per brand rules
logo_text(sl, 0.55, 0.22, size=22, color=WHITE)
# REMOVED hidden placeholder

# Actual logo (full colour) — place on dark, will show in lilac tones
logo(sl, 0.5, 0.18, 3.0)

txt(sl, 'INVESTOR PRESENTATION  ·  2026',
    0.55, 1.42, 7, 0.32,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)

# Hero headline — Exo 2, white + lilac split
txt(sl, 'The trusted path to', 0.55, 1.9, 7, 0.85,
    size=40, bold=True, color=WHITE, font=EXO2)
txt(sl, 'Asian beauty discovery.', 0.55, 2.72, 7.2, 0.85,
    size=40, bold=True, color=LILAC, font=EXO2)

txt(sl,
    'Simplex-ity connects verified creators with beauty brands\n'
    'through curated campaigns, TINT virtual try-on,\n'
    'and a reliability-scored influencer ecosystem.',
    0.55, 3.75, 7.0, 1.4,
    size=14.5, color=SOFT_LILAC, font=MONTSERRAT)

bar(sl, 0.55, 5.35, 6.8, color=LILAC)

txt(sl,
    'Kieran  ·  kieran@5senses.global\n'
    'SIMPLEX-ITY by 5SENSESBEAUTY LIMITED  ·  Hong Kong',
    0.55, 5.52, 7, 0.7,
    size=12, color=MUTED, font=MONTSERRAT)
txt(sl, 'simplex-ity.fluentlab.co',
    0.55, 6.28, 5, 0.38, size=12, color=SOFT_LILAC, font=MONTSERRAT)

# Right panel — brand tagline
txt(sl, 'Asian beauty, made easier\nto trust, learn, and love.',
    8.1, 2.8, 5.0, 1.6,
    size=20, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font=EXO2)
bar(sl, 8.5, 4.6, 4.1, color=LILAC)
txt(sl, 'SIMPLIFY TO AMPLIFY',
    8.1, 4.75, 5.0, 0.38,
    size=12, bold=True, color=LILAC,
    align=PP_ALIGN.CENTER, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, 'EXECUTIVE SUMMARY')
headline(sl, 'The opportunity in one page.')
sub(sl, 'Three markets. One platform. Built for trust.')

rect(sl, 0.3, 1.58, 12.73, 0.95, LAV_WASH)
bar(sl, 0.3, 1.58, 0.06, h=0.95, color=VIOLET)
txt(sl,
    'Asian beauty, made easier to trust, learn, and love.',
    0.52, 1.72, 12.3, 0.68,
    size=20, bold=True, color=VIOLET, font=EXO2)

boxes = [
    ('The Problem',
     'Brands cannot verify influencer ROI. Consumers cannot trust shade fit. No curated, guided platform connects both.'),
    ('Our Solution',
     'Simplex-ity is the managed influencer marketplace for Asian beauty — creator matching, TINT try-on, and real-time campaign tracking.'),
    ('The Business Model',
     'Brands pay HK$400-$2,500/month membership. Influencers earn 10-15% commission. Platform retains 5% performance fee post-trial.'),
    ('The Market',
     '$2.3B AI beauty market. $2.03B influencer platform market by 2031. HK Instagram at 45.3% share. TikTok Shop beauty up 108%.'),
]
for i, (title, body) in enumerate(boxes):
    x = 0.3 if i % 2 == 0 else 6.82
    y = 2.75 if i < 2 else 5.0
    feature_card(sl, title, body, x, y, w=6.2, h=2.05)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '01 — THE PROBLEM')
headline(sl, 'A $2.3B market without a trusted guide.')
sub(sl, 'Three groups. Three unresolved problems. One platform to fix all three.')

cols = [
    ('For Brands', VIOLET, [
        'No reliable way to verify influencer audience quality.',
        'Virtual try-on is expensive and disconnected from campaigns.',
        'Campaign ROI is measured manually, weeks after the fact.',
        'Influencer reliability tracked in spreadsheets.',
    ]),
    ('For Influencers', LILAC, [
        'Commission platforms offer no tech advantage.',
        'Brand matching is opaque and relationship-dependent.',
        'No tools to demonstrate professional reliability.',
        'Income inconsistent with no retainer structure.',
    ]),
    ('For Consumers', SOFT_LILAC, [
        'Shade uncertainty stops purchase at the final step.',
        'Creator content is aspirational but not always verified.',
        'Asian beauty discovery is fragmented across platforms.',
        'No curated source guiding from discovery to purchase.',
    ]),
]
x_pos = [0.3, 4.65, 9.0]
for i, (group, accent, points) in enumerate(cols):
    x = x_pos[i]
    rect(sl, x, 2.08, 4.1, 5.15, LAV_WASH)
    bar(sl, x, 2.08, 4.1, color=accent)
    txt(sl, group, x+0.18, 2.22, 3.72, 0.48,
        size=17, bold=True, color=accent, font=EXO2)
    y = 2.86
    for pt in points:
        txt(sl, '— ' + pt, x+0.18, y, 3.7, 0.58,
            size=12, color=BODY, font=MONTSERRAT)
        y += 0.64

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '02 — THE SOLUTION')
headline(sl, 'Simplex-ity: one managed platform.')
sub(sl, 'Brands. Influencers. Try-on technology. Unified.')

pillars = [
    ('Verified Creator Matching',
     'Curated influencers matched by audience data and platform reliability score. Selected, not scraped.'),
    ('TINT Virtual Try-On',
     'AR makeup try-on embedded in creator content and brand pages. 25%+ higher conversion rate vs standard links.'),
    ('Reliability Scorecard',
     'Every creator rated 0-100% on delivery and conversion. Below 70% triggers automatic shop suspension.'),
    ('End-to-End Tracking',
     'Live counts, blog delivery, commission paid, and full post-campaign ROI in one real-time dashboard.'),
    ('Brand Membership Tiers',
     'Structured brand partnerships from HK$400/month (Essential) to HK$2,500/month (Enterprise).'),
    ('Digital Franchise Model',
     'Influencers earn 10-15% commission with zero upfront cost. Platform grows through creator cross-pollination.'),
]
for i, (title, body) in enumerate(pillars):
    x = 0.3 + (i % 3) * 4.35
    y = 2.08 if i < 3 else 4.72
    rect(sl, x, y, 4.1, 2.42, LAV_WASH)
    bar(sl, x, y, 4.1)
    txt(sl, title, x+0.18, y+0.12, 3.72, 0.5,
        size=13.5, bold=True, color=VIOLET, font=EXO2)
    txt(sl, body, x+0.18, y+0.68, 3.72, 1.65,
        size=12, color=BODY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — TINT TECHNOLOGY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
bar(sl, 0, 0, 13.33)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '03 — KEY DIFFERENTIATOR')
headline(sl, 'TINT Virtual Try-On.', color=BODY, font=EXO2)
sub(sl, 'The only Asian beauty platform with embedded AR makeup technology.')

stats = [
    ('25%+', 'Higher conversion\nvs standard links'),
    ('200%', 'More conversions\nvs static images'),
    ('30%',  'Higher average\norder value'),
    ('17%',  'AI beauty market\nCAGR 2025-2033'),
]
for i, (val, label) in enumerate(stats):
    stat_card(sl, val, label, 0.3 + i*3.25, 2.05,
              bg_col=WHITE)

rect(sl, 0.3, 4.42, 12.73, 2.82, WHITE)
bar(sl, 0.3, 4.42, 12.73, color=LILAC)
txt(sl, 'What TINT delivers for your brand partners',
    0.5, 4.55, 12, 0.42,
    size=14, bold=True, color=BODY, font=EXO2)

tint_items = [
    ('Real-Time AR Try-On',
     'Lipstick, eyeshadow, foundation try-on inside creator posts. No app download needed.'),
    ('AI Skin Analysis',
     'Personalised shade recommendations at point of discovery, based on consumer skin tone.'),
    ('Embedded Commerce',
     'Try-on links directly to purchase. Removes friction between inspiration and conversion.'),
    ('Creator Integration',
     'Runs live inside streams and blog posts. Creators use it live; consumers trust what they see.'),
]
for i, (title, desc) in enumerate(tint_items):
    x = 0.5 + i*3.2
    txt(sl, title, x, 5.1, 3.0, 0.42,
        size=12.5, bold=True, color=VIOLET, font=EXO2)
    txt(sl, desc, x, 5.58, 3.0, 1.4,
        size=11.5, color=BODY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '04 — MARKET OPPORTUNITY')
headline(sl, 'A market ready to be shaped.')
sub(sl, 'Sized. Validated. Underserved.')

mkt = [
    ('$2.3B',  'Global AI Beauty\nPersonalisation\nMarket 2026'),
    ('$2.03B', 'Influencer Platform\nMarket by 2031'),
    ('108%',   'TikTok Shop beauty\nsales growth 2025'),
    ('68%',    'Beauty purchases\ninfluenced by social\ndiscovery'),
]
for i, (val, label) in enumerate(mkt):
    stat_card(sl, val, label, 0.3 + i*3.25, 1.68)

rect(sl, 0.3, 3.98, 12.73, 1.08, LAV_WASH)
bar(sl, 0.3, 3.98, 12.73, color=LILAC)
txt(sl, 'Hong Kong Market Context', 0.5, 4.1, 12, 0.35,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)
txt(sl,
    'Instagram leads HK influencer marketing at 45.3% share. '
    'Xiaohongshu is the fastest-growing platform in HK beauty. '
    'Simplex-ity operates across both. 68% of global beauty purchases '
    'are influenced at the discovery stage — the moment Simplex-ity captures.',
    0.5, 4.5, 12.3, 0.5,
    size=12.5, color=BODY, font=MONTSERRAT)

rect(sl, 0.3, 5.22, 12.73, 2.03, WHITE)
bar(sl, 0.3, 5.22, 12.73, color=NEUTRAL)
txt(sl, 'Why now: the gap in the market',
    0.5, 5.35, 12, 0.38,
    size=13, bold=True, color=BODY, font=EXO2)
gaps = [
    'Standard affiliate platforms offer no AR technology or reliability governance.',
    'Beauty marketplaces are catalogue-driven and not creator-guided.',
    'No platform combines TINT try-on with influencer reliability scoring for the Asian beauty category.',
]
y = 5.82
for g in gaps:
    txt(sl, '— ' + g, 0.5, y, 12.3, 0.36,
        size=12.5, color=SECONDARY, font=MONTSERRAT)
    y += 0.42

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '05 — BUSINESS MODEL')
headline(sl, 'Three revenue streams. Zero upfront risk.')

streams = [
    ('Stream 1: Brand Membership (Recurring SaaS)', VIOLET,
     'Essential HK$400/mo   |   Professional HK$900/mo   |   Enterprise HK$2,500/mo',
     'Recurring monthly subscription. Predictable, scalable, platform-independent.'),
    ('Stream 2: Performance Fee (Commission Take-Rate)', LILAC,
     '5% platform fee retained from influencer commissions after 3-month trial. Influencer commissions range 10-15% of net sales.',
     'Scales automatically with GMV. No sales team required to grow this line.'),
    ('Stream 3: Premium Shop Subscriptions (Creator Upsell)', SOFT_LILAC,
     'HK$20/month per influencer for advanced analytics, custom storefront UI, and deep-dive campaign data.',
     'Optional upsell unlocked once creators are earning. Low friction, high margin.'),
]
y = 1.55
for title, accent, detail, note in streams:
    rect(sl, 0.3, y, 12.73, 1.62, LAV_WASH)
    bar(sl, 0.3, y, 0.06, h=1.62, color=accent)
    txt(sl, title, 0.52, y+0.1, 10, 0.44,
        size=14, bold=True, color=VIOLET, font=EXO2)
    txt(sl, detail, 0.52, y+0.6, 12.3, 0.42,
        size=13, color=BODY, font=MONTSERRAT)
    txt(sl, note, 0.52, y+1.08, 12.3, 0.42,
        size=12, color=SECONDARY, italic=True, font=MONTSERRAT)
    y += 1.82

rect(sl, 0.3, 7.0, 12.73, 0.38, LAV_WASH)
txt(sl,
    'All three streams compound: more brands attract more creators, '
    'which grow GMV, which increases performance fees.',
    0.5, 7.05, 12.5, 0.3,
    size=12, bold=True, color=BODY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
bar(sl, 0, 0, 13.33)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '06 — HOW IT WORKS')
headline(sl, 'From brief to results in four steps.')

steps = [
    ('01', 'Brand Brief\nand Matching',
     'Brand shares product, budget, and target consumer. Platform matches to verified creators using reliability score and audience data.'),
    ('02', 'TINT Try-On\nSetup',
     'TINT virtual try-on configured for the product and embedded in creator content. Consumers experience the product before buying.'),
    ('03', 'Campaign\nExecution',
     'Creators deliver 1 review blog and 2 professional live streams per campaign window. Shared calendar governs all deadlines.'),
    ('04', 'Performance\nReport',
     'Full post-campaign ROI: reach, conversions, commission paid, reliability score impact, and next-campaign recommendations.'),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.3 + i * 3.25
    rect(sl, x, 1.85, 3.0, 5.35, WHITE)
    bar(sl, x, 1.85, 3.0)
    txt(sl, num, x+0.15, 2.05, 2.7, 0.75,
        size=36, bold=True, color=LILAC,
        align=PP_ALIGN.CENTER, font=EXO2)
    txt(sl, title, x+0.15, 2.88, 2.7, 0.75,
        size=15, bold=True, color=BODY, font=EXO2)
    txt(sl, desc, x+0.15, 3.72, 2.7, 2.85,
        size=12, color=SECONDARY, font=MONTSERRAT)
    if i < 3:
        txt(sl, '→', x+3.05, 4.3, 0.42, 0.52,
            size=22, color=LILAC, align=PP_ALIGN.CENTER,
            font=MONTSERRAT)

bar(sl, 0, 7.22, 13.33, h=0.28, color=VIOLET)
txt(sl,
    'Content Trinity: 1x Detailed Blog   +   2x Professional Live Streams   +   AI Sales Certification   +   Reliability Score 0-100%',
    0.3, 7.24, 12.73, 0.24,
    size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — TIER STRUCTURE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '07 — PLATFORM TIERS')
headline(sl, 'Two ecosystems. Both governed.')

txt(sl, 'BRAND MEMBERSHIP', 0.3, 1.52, 6, 0.32,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)

brand_tiers = [
    ('Essential',    'HK$400/mo',
     ['1 campaign/quarter', 'Up to 5 creators', 'TINT basic']),
    ('Professional', 'HK$900/mo',
     ['2 campaigns/quarter', 'Up to 15 creators', 'Full TINT + dashboard']),
    ('Enterprise',   'HK$2,500/mo',
     ['Unlimited campaigns', 'Unlimited creators', 'Dedicated manager']),
]
x_pos = [0.3, 4.55, 8.8]
for i, (name, price, feats) in enumerate(brand_tiers):
    x = x_pos[i]
    featured = (i == 1)
    bc = SOFT_LILAC if featured else LAV_WASH
    ac = VIOLET
    rect(sl, x, 1.9, 4.0, 2.58, bc)
    bar(sl, x, 1.9, 4.0, color=ac)
    txt(sl, name, x+0.15, 2.04, 3.7, 0.48,
        size=18, bold=True, color=ac,
        align=PP_ALIGN.CENTER, font=EXO2)
    rect(sl, x+0.15, 2.6, 3.7, 0.42, VIOLET)
    txt(sl, price, x+0.15, 2.62, 3.7, 0.38,
        size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, font=MONTSERRAT)
    for j, feat in enumerate(feats):
        txt(sl, '   ' + feat, x+0.18,
            3.1 + j*0.42, 3.7, 0.38,
            size=12, color=BODY, font=MONTSERRAT)

txt(sl, 'INFLUENCER TIERS', 0.3, 4.68, 6, 0.32,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)

inf_tiers = [
    ('Rising Star',     'Essential brands',
     ['1 Blog + 1 Live/month', '10% commission']),
    ('Star Influencer', 'Professional brands',
     ['2 Blogs + 4 Lives/month', '12% + content fees']),
    ('Superstar',       'Enterprise brands',
     ['5 Blogs + 10 Lives/month', '15% + 1st Launch access']),
]
for i, (name, brand_match, feats) in enumerate(inf_tiers):
    x = x_pos[i]
    rect(sl, x, 5.08, 4.0, 2.12, LAV_WASH)
    bar(sl, x, 5.08, 4.0, color=LILAC)
    txt(sl, name, x+0.15, 5.2, 3.7, 0.45,
        size=15, bold=True, color=VIOLET, font=EXO2)
    txt(sl, brand_match, x+0.18, 5.72, 3.7, 0.35,
        size=12, color=SECONDARY, italic=True, font=MONTSERRAT)
    for j, feat in enumerate(feats):
        txt(sl, '   ' + feat, x+0.18,
            6.12 + j*0.42, 3.7, 0.38,
            size=12, color=BODY, font=MONTSERRAT)

bar(sl, 0, 7.27, 13.33, h=0.23, color=VIOLET)
txt(sl,
    'Red Card rule: reliability score below 70% = permanent shop suspension.',
    0.3, 7.29, 12.73, 0.2,
    size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITIVE MOAT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
bar(sl, 0, 0, 13.33)
bar(sl, 0, 0.07, 0.06, h=7.43, color=LILAC)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '08 — COMPETITIVE ADVANTAGE')
headline(sl, 'Five moats. Built in from day one.')
sub(sl, 'Curated by design. Not replicable by horizontal platforms.')

moats = [
    ('Try-On Technology',
     'TINT virtual AR makeup integrated at platform level. Competitors must build or license this separately.'),
    ('Reliability Governance',
     'The 0-100% scorecard is proprietary. No other Asian beauty platform governs creator behaviour this systematically.'),
    ('Category Focus',
     'Built exclusively for Asian beauty. Category depth creates trust that horizontal affiliate platforms cannot replicate.'),
    ('Network Effects',
     'Cross-pollination drives follower growth for every creator. The more creators on platform, the more valuable each one becomes.'),
    ('Data Advantage',
     'Every campaign generates proprietary data on what converts in Asian beauty by product type, skin tone, and creator tier.'),
]
for i, (title, desc) in enumerate(moats):
    y = 1.65 + i*1.05
    rect(sl, 0.3, y, 12.73, 0.97, LAV_WASH)
    bar(sl, 0.3, y, 0.06, h=0.97, color=VIOLET)
    txt(sl, title, 0.52, y+0.1, 3.5, 0.42,
        size=14, bold=True, color=VIOLET, font=EXO2)
    txt(sl, desc, 4.1, y+0.12, 8.8, 0.7,
        size=12.5, color=BODY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — ROADMAP
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
bar(sl, 0, 0, 13.33)
logo(sl, 0.3, 0.15, 1.9)
eyebrow(sl, '09 — ROADMAP')
headline(sl, 'Where we are. Where we are going.')

roadmap = [
    ('Q2 2026\nNow', VIOLET,
     'Platform live. Brand partnerships open. First influencer cohort onboarding. TINT integration active.'),
    ('Q3 2026', LILAC,
     'First 10 brand members signed. 50 verified creators active. Campaign dashboard launched. First ROI reports delivered.'),
    ('Q4 2026', LILAC,
     'Professional and Enterprise brands onboarded. Retainer model active for top creators. XHS and TikTok channels added.'),
    ('2027', SOFT_LILAC,
     'Regional expansion: Singapore and Taiwan. Institutional brand partnerships. Series A readiness achieved.'),
]
for i, (phase, accent, desc) in enumerate(roadmap):
    x = 0.3 + i*3.25
    rect(sl, x, 1.82, 3.0, 5.42, WHITE)
    bar(sl, x, 1.82, 3.0, color=accent)
    col = VIOLET if i == 0 else BODY
    txt(sl, phase, x+0.15, 2.02, 2.7, 0.75,
        size=18, bold=True, color=col, font=EXO2)
    if i == 0:
        rect(sl, x+0.15, 2.85, 1.8, 0.3, VIOLET)
        txt(sl, 'ACTIVE NOW', x+0.15, 2.87, 1.8, 0.26,
            size=9, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=MONTSERRAT)
    txt(sl, desc, x+0.15, 3.28, 2.7, 3.8,
        size=12.5, color=SECONDARY, font=MONTSERRAT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — THE ASK / CLOSING
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK_HERO)
bar(sl, 0, 0, 13.33, color=VIOLET)
rect(sl, 7.5, 0, 5.83, 7.5, DARK2)
bar(sl, 7.5, 0, 0.06, h=7.5, color=LILAC)

logo(sl, 0.5, 0.2, 2.8)

txt(sl, 'THE INVESTMENT OPPORTUNITY',
    0.55, 1.38, 6.7, 0.35,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)

txt(sl, 'Join Simplex-ity at\nthe foundation stage.',
    0.55, 1.82, 6.7, 1.75,
    size=38, bold=True, color=WHITE, font=EXO2)

txt(sl,
    'We are seeking seed investment to accelerate brand onboarding, '
    'expand the creator ecosystem across HK, Singapore, and Taiwan, '
    'and build the full analytics infrastructure.',
    0.55, 3.7, 6.7, 1.2,
    size=14, color=SOFT_LILAC, font=MONTSERRAT)

bar(sl, 0.55, 5.08, 6.5, color=LILAC)

use_items = [
    ('Brand Sales:', 'Onboard first 20 brand partners across all tiers'),
    ('Creator Growth:', 'Build verified creator cohort to 200 active influencers'),
    ('Technology:', 'Full TINT integration, analytics dashboard, reporting suite'),
    ('Regional:', 'Singapore and Taiwan market entry in 2027'),
]
y = 5.28
for label, desc in use_items:
    txt(sl, label, 0.55, y, 2.4, 0.35,
        size=12, bold=True, color=LILAC, font=EXO2)
    txt(sl, desc, 3.05, y, 4.2, 0.35,
        size=12, color=SOFT_LILAC, font=MONTSERRAT)
    y += 0.42

# Right panel
txt(sl, 'KEY NUMBERS', 7.7, 1.35, 5.3, 0.32,
    size=10, bold=True, color=MUTED, font=MONTSERRAT)

key = [
    ('3 Revenue Streams', 'Membership, Performance Fee, Premium Shop'),
    ('HK$400 to $2,500', 'Brand membership per month'),
    ('10-15%', 'Influencer commission rate'),
    ('25%+', 'TINT conversion rate uplift'),
    ('$0', 'Upfront cost for influencers'),
    ('$2.3B', 'Target market size 2026'),
]
y = 1.82
for label, val in key:
    rect(sl, 7.68, y, 5.3, 0.75, RGBColor(0x28, 0x18, 0x48))
    bar(sl, 7.68, y, 0.05, h=0.75, color=LILAC)
    txt(sl, label, 7.86, y+0.05, 3.8, 0.35,
        size=14, bold=True, color=WHITE, font=EXO2)
    txt(sl, val, 7.86, y+0.42, 4.9, 0.28,
        size=11.5, color=SOFT_LILAC, font=MONTSERRAT)
    y += 0.85

txt(sl, 'kieran@5senses.global',
    0.55, 6.98, 6.7, 0.32, size=13, color=SOFT_LILAC, font=MONTSERRAT)
txt(sl, 'simplex-ity.fluentlab.co   |   Kwun Tong, Hong Kong',
    0.55, 7.22, 6.7, 0.25, size=11, color=MUTED, font=MONTSERRAT)

prs.save('/app/SIMPLEX-ITY_Investor_Pitch_v2.pptx')
print('Done! 12 slides saved.')
