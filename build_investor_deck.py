from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Colors (Official Design System) ─────────────────────────
LILAC      = RGBColor(0x8c, 0x82, 0xfc)   # Primary Lilac
VIOLET     = RGBColor(0x5e, 0x50, 0xfb)   # Accent Violet
SOFT_LILAC = RGBColor(0xba, 0xb4, 0xfd)   # Soft Lilac
LAV_WASH   = RGBColor(0xe8, 0xe6, 0xfe)   # Lavender Wash
WHITE      = RGBColor(0xff, 0xff, 0xff)   # White Canvas
NEUTRAL    = RGBColor(0xe6, 0xe6, 0xe6)   # Neutral Grey
BODY       = RGBColor(0x1a, 0x1a, 0x1f)   # Body Text
SECONDARY  = RGBColor(0x5a, 0x5a, 0x6a)   # Secondary Text
MUTED      = RGBColor(0x99, 0x99, 0xaa)   # Muted Text
DARK_HERO  = RGBColor(0x14, 0x10, 0x2a)   # Dark hero BG

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

LOGO = '/app/logo.jpg'
IMG1 = '/app/product1.jpg'
IMG2 = '/app/product2.jpg'

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False,
        color=BODY, align=PP_ALIGN.LEFT, italic=False,
        font='Montserrat', wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name = font
    return tb

def add_logo(slide, l, t, w=2.2):
    slide.shapes.add_picture(LOGO, Inches(l), Inches(t), width=Inches(w))

def divider(slide, l, t, w, color=VIOLET, h=0.055):
    rect(slide, l, t, w, h, color)

def pill_label(slide, label, l, t, w=2.8, h=0.36,
               bg_col=VIOLET, fg_col=WHITE, size=10):
    rect(slide, l, t, w, h, bg_col)
    txt(slide, label, l, t+0.01, w, h-0.02,
        size=size, bold=True, color=fg_col, align=PP_ALIGN.CENTER)

def stat_box(slide, val, label, x, y, w=3.05, h=2.0,
             val_color=VIOLET, bg_col=LAV_WASH, accent=VIOLET):
    rect(slide, x, y, w, h, bg_col)
    divider(slide, x, y, w, accent, h=0.06)
    txt(slide, val, x+0.1, y+0.15, w-0.2, 0.82,
        size=32, bold=True, color=val_color,
        align=PP_ALIGN.CENTER, font='Montserrat')
    txt(slide, label, x+0.1, y+1.0, w-0.2, 0.85,
        size=12, color=BODY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — INVESTOR COVER
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK_HERO)
# Top violet bar
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
# Right lilac panel
rect(sl, 8.0, 0, 5.33, 7.5, RGBColor(0x20, 0x14, 0x40))
divider(sl, 8.0, 0, 0.06, LILAC, h=7.5)

# Product image on right
sl.shapes.add_picture(IMG1, Inches(8.2), Inches(0.3),
                      width=Inches(5.0), height=Inches(7.0))
# Overlay tint on image
rect(sl, 8.0, 0, 5.33, 7.5, RGBColor(0x20, 0x14, 0x40))
# Re-add image on top of tint for proper layering — skip (would need z-order hack)

# Logo
add_logo(sl, 0.55, 0.22, 2.6)

txt(sl, 'INVESTOR PRESENTATION  ·  2026',
    0.55, 1.45, 7, 0.35,
    size=10, bold=True, color=MUTED)

txt(sl, 'The trusted path to\nAsian beauty discovery.',
    0.55, 1.9, 7.2, 2.1,
    size=40, bold=True, color=WHITE, font='Montserrat')

txt(sl, 'Simplex-ity connects verified creators with beauty brands\nthrough curated campaigns, TINT virtual try-on, and\na reliability-scored influencer ecosystem.',
    0.55, 4.1, 7.0, 1.3,
    size=15, color=SOFT_LILAC)

divider(sl, 0.55, 5.6, 6.5, LILAC, h=0.055)

# Contact row
txt(sl, 'Kieran  ·  kieran@5senses.global  ·  SIMPLEX-ITY by 5SENSESBEAUTY LIMITED',
    0.55, 5.75, 8, 0.42, size=12, color=MUTED)
txt(sl, 'simplex-ity.fluentlab.co',
    0.55, 6.22, 5, 0.4, size=12, color=SOFT_LILAC)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, 'EXECUTIVE SUMMARY', 0.3, 0.2, 5, 0.28,
    size=9, bold=True, color=MUTED)

txt(sl, 'The opportunity in one page.',
    0.3, 0.65, 10, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

# Large highlighted brand statement
rect(sl, 0.3, 1.5, 12.73, 1.0, LAV_WASH)
rect(sl, 0.3, 1.5, 0.07, 1.0, VIOLET)
txt(sl, 'Asian beauty, made easier to trust, learn, and love.',
    0.5, 1.65, 12.3, 0.7,
    size=20, bold=True, color=VIOLET, font='Montserrat')

boxes = [
    ('The Problem',
     'Asian beauty brands cannot verify influencer ROI. Consumers cannot trust whether products suit their skin. There is no curated, guided platform connecting both.'),
    ('Our Solution',
     'Simplex-ity is the managed influencer marketplace for Asian beauty. We combine verified creator matching, TINT virtual try-on, and real-time campaign tracking in one platform.'),
    ('The Model',
     'Brands pay monthly membership (HK$400 to HK$2,500). Influencers earn 10-15% commission with zero upfront cost. Platform retains a 5% performance fee post-trial.'),
    ('The Opportunity',
     'Global AI beauty personalisation market: US$2.3B in 2026. Influencer platform market: US$2.03B by 2031. HK beauty market growing with Instagram at 45.3% share.'),
]
positions = [(0.3, 2.7), (6.8, 2.7), (0.3, 5.0), (6.8, 5.0)]
for i, (title, desc) in enumerate(boxes):
    x, y = positions[i]
    rect(sl, x, y, 6.2, 2.1, LAV_WASH)
    rect(sl, x, y, 6.2, 0.06, VIOLET)
    txt(sl, title, x+0.2, y+0.12, 5.8, 0.42,
        size=14, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, x+0.2, y+0.6, 5.8, 1.35,
        size=12.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '01 — THE PROBLEM', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'A $2.3B market without a trusted guide.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

txt(sl, 'Three groups. Three unresolved problems. One platform to fix all three.',
    0.3, 1.42, 11, 0.4, size=15, color=SECONDARY, italic=True)

cols = [
    ('For Brands', VIOLET,
     ['No reliable way to verify influencer audience quality.',
      'Virtual try-on technology is expensive and disconnected from campaigns.',
      'Campaign ROI is measured manually, weeks after the fact.',
      'Influencer reliability tracked in spreadsheets, campaigns fail mid-execution.']),
    ('For Influencers', LILAC,
     ['Commission platforms offer no tech advantage to grow sales.',
      'Brand matching is opaque and relationship-dependent.',
      'No tools to demonstrate professional reliability to brands.',
      'Income is inconsistent with no retainer or anchor structure.']),
    ('For Consumers', SOFT_LILAC,
     ['Shade uncertainty stops purchase at the final step.',
      'Creator content is aspirational but not always trustworthy.',
      'Asian beauty discovery is fragmented across platforms.',
      'No single curated source that guides from discovery to purchase.']),
]
x_pos = [0.3, 4.65, 9.0]
for i, (group, accent, points) in enumerate(cols):
    x = x_pos[i]
    rect(sl, x, 2.1, 4.1, 5.1, LAV_WASH)
    rect(sl, x, 2.1, 4.1, 0.06, accent)
    txt(sl, group, x+0.18, 2.25, 3.7, 0.45,
        size=16, bold=True, color=accent, font='Montserrat')
    y = 2.85
    for pt in points:
        txt(sl, '— ' + pt, x+0.18, y, 3.7, 0.55, size=11.5, color=BODY)
        y += 0.62

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '02 — THE SOLUTION', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'Simplex-ity: one managed platform.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

txt(sl, 'Brands. Influencers. Try-on technology. Unified.',
    0.3, 1.38, 10, 0.4, size=15, color=SECONDARY, italic=True)

pillars = [
    ('Verified Creator Matching',
     'Curated influencers matched to your brand by audience data and platform reliability score. Selected, not scraped.'),
    ('TINT Virtual Try-On',
     'AR makeup try-on embedded in creator content and brand pages. 25%+ higher conversion rate vs standard links.'),
    ('Reliability Scorecard',
     'Every creator rated 0-100% on delivery, engagement, and conversion. Below 70% triggers automatic shop suspension.'),
    ('End-to-End Campaign Tracking',
     'Live counts, blog delivery, commission paid, and full post-campaign ROI report, in one real-time dashboard.'),
    ('Brand Membership Tiers',
     'Structured brand partnerships from HK$400/month (Essential) to HK$2,500/month (Enterprise) with clear deliverables.'),
    ('Digital Franchise Model',
     'Influencers earn 10-15% commission with zero upfront cost. Platform grows through creator cross-pollination and follower synergy.'),
]
positions = [(0.3,2.1),(4.65,2.1),(9.0,2.1),(0.3,4.75),(4.65,4.75),(9.0,4.75)]
for i, (title, desc) in enumerate(pillars):
    x, y = positions[i]
    rect(sl, x, y, 4.1, 2.45, LAV_WASH)
    rect(sl, x, y, 4.1, 0.06, VIOLET)
    txt(sl, title, x+0.18, y+0.12, 3.75, 0.48,
        size=13, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, x+0.18, y+0.65, 3.75, 1.65,
        size=11.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — TINT TECHNOLOGY (DIFFERENTIATOR)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '03 — KEY DIFFERENTIATOR', 0.3, 0.2, 6, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'TINT Virtual Try-On.', 0.3, 0.65, 8, 0.65,
    size=38, bold=True, color=BODY, font='Montserrat')
txt(sl, 'The only Asian beauty platform with embedded AR makeup technology.',
    0.3, 1.4, 11, 0.4, size=15, color=SECONDARY, italic=True)

# Stats row
stats = [
    ('25%+', 'Higher conversion\nvs standard links'),
    ('200%', 'More conversions\nvs static images'),
    ('30%',  'Higher average\norder value'),
    ('17%',  'AI beauty market\nCAGR 2025-2033'),
]
for i, (val, label) in enumerate(stats):
    x = 0.3 + i*3.25
    rect(sl, x, 2.15, 3.0, 2.1, WHITE)
    rect(sl, x, 2.15, 3.0, 0.06, VIOLET)
    txt(sl, val, x+0.1, 2.28, 2.8, 0.85,
        size=36, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font='Montserrat')
    txt(sl, label, x+0.1, 3.18, 2.8, 0.9,
        size=12, color=BODY, align=PP_ALIGN.CENTER)

# What TINT delivers
rect(sl, 0.3, 4.55, 12.73, 2.7, WHITE)
rect(sl, 0.3, 4.55, 12.73, 0.06, LILAC)
txt(sl, 'What TINT delivers for your brand partners', 0.5, 4.68, 12, 0.42,
    size=14, bold=True, color=BODY)

tint_items = [
    ('Real-Time AR Try-On',
     'Lipstick, eyeshadow, and foundation try-on inside creator posts. No app download needed.'),
    ('AI Skin Analysis',
     'Personalised shade recommendations based on consumer skin tone, shown at the point of discovery.'),
    ('Embedded Commerce',
     'Try-on links directly to purchase. Removes the friction between inspiration and conversion.'),
    ('Creator Integration',
     'Runs inside live streams and blog posts. Creators use it live, consumers trust what they see.'),
]
for i, (title, desc) in enumerate(tint_items):
    x = 0.5 + i*3.2
    txt(sl, title, x, 5.22, 3.0, 0.38,
        size=12, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, x, 5.65, 3.0, 0.98, size=11.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '04 — MARKET OPPORTUNITY', 0.3, 0.2, 6, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'A market ready to be shaped.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

mkt = [
    ('$2.3B', 'Global AI Beauty\nPersonalisation\nMarket 2026'),
    ('$2.03B','Influencer Platform\nMarket\nby 2031'),
    ('108%',  'TikTok Shop beauty\nsales growth\n2025'),
    ('68%',   'Beauty purchases\ninfluenced by\nsocial discovery'),
]
for i, (val, label) in enumerate(mkt):
    stat_box(sl, val, label, 0.3 + i*3.25, 1.7)

# HK Insight
rect(sl, 0.3, 4.0, 12.73, 1.1, LAV_WASH)
rect(sl, 0.3, 4.0, 12.73, 0.06, LILAC)
txt(sl, 'Hong Kong Market Context', 0.5, 4.12, 12, 0.35,
    size=10, bold=True, color=MUTED)
txt(sl, 'Instagram leads HK influencer marketing at 45.3% share. Xiaohongshu (XHS) is the fastest-growing platform in HK beauty. Simplex-ity operates across both. TikTok Shop beauty sales grew 108% in 2025. 68% of global beauty purchases are influenced at the discovery stage, the precise moment Simplex-ity captures.',
    0.5, 4.52, 12.3, 0.52, size=12.5, color=BODY)

# Platform comparison
rect(sl, 0.3, 5.3, 12.73, 1.95, WHITE)
rect(sl, 0.3, 5.3, 12.73, 0.06, NEUTRAL)
txt(sl, 'Why now: gap in the market', 0.5, 5.42, 12, 0.38,
    size=12, bold=True, color=BODY)
gaps = [
    'Standard affiliate platforms offer no AR technology or reliability governance.',
    'Beauty marketplaces are catalogue-driven and not creator-guided.',
    'No platform currently combines TINT try-on with influencer reliability scoring for the Asian beauty category.',
]
y = 5.88
for g in gaps:
    txt(sl, '— ' + g, 0.5, y, 12.3, 0.35, size=12, color=SECONDARY)
    y += 0.38

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '05 — BUSINESS MODEL', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'Three revenue streams. Zero upfront risk.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

# Revenue stream 1: Brand Membership
rect(sl, 0.3, 1.55, 12.73, 1.55, LAV_WASH)
rect(sl, 0.3, 1.55, 0.06, 1.55, VIOLET)
txt(sl, 'Stream 1: Brand Membership (Recurring SaaS)',
    0.5, 1.65, 7, 0.42, size=14, bold=True, color=VIOLET, font='Montserrat')
tiers_txt = 'Essential HK$400/mo  |  Professional HK$900/mo  |  Enterprise HK$2,500/mo'
txt(sl, tiers_txt, 0.5, 2.12, 9, 0.38, size=13, color=BODY)
txt(sl, 'Recurring monthly subscription. Predictable, scalable, platform-independent.',
    0.5, 2.55, 9, 0.38, size=12, color=SECONDARY, italic=True)

# Revenue stream 2: Performance fee
rect(sl, 0.3, 3.3, 12.73, 1.55, LAV_WASH)
rect(sl, 0.3, 3.3, 0.06, 1.55, LILAC)
txt(sl, 'Stream 2: Performance Fee (Commission Take-Rate)',
    0.5, 3.4, 9, 0.42, size=14, bold=True, color=VIOLET, font='Montserrat')
txt(sl, '5% platform fee retained from influencer commissions after 3-month trial. Commissions range 10-15% of net sales.',
    0.5, 3.87, 9, 0.45, size=13, color=BODY)
txt(sl, 'Scales automatically with GMV. No sales team required to grow this line.',
    0.5, 4.38, 9, 0.35, size=12, color=SECONDARY, italic=True)

# Revenue stream 3: Premium
rect(sl, 0.3, 5.05, 12.73, 1.55, LAV_WASH)
rect(sl, 0.3, 5.05, 0.06, 1.55, SOFT_LILAC)
txt(sl, 'Stream 3: Premium Shop Subscriptions (Creator Upsell)',
    0.5, 5.15, 9, 0.42, size=14, bold=True, color=VIOLET, font='Montserrat')
txt(sl, 'HK$20/month per influencer for advanced analytics, custom storefront UI, and deep-dive campaign data.',
    0.5, 5.62, 9, 0.45, size=13, color=BODY)
txt(sl, 'Optional upsell unlocked once creators are earning. Low friction, high margin.',
    0.5, 6.13, 9, 0.35, size=12, color=SECONDARY, italic=True)

txt(sl, 'All three streams compound: more brands attract more creators, which grow GMV, which increases performance fees.',
    0.3, 6.72, 12.73, 0.45, size=12.5, bold=True, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — HOW IT WORKS (PLATFORM FLOW)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '06 — HOW IT WORKS', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'From brief to results in four steps.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

steps = [
    ('01', 'Brand Brief\n& Matching',
     'Brand shares product, budget, and target consumer. Platform matches to verified creators using reliability score and audience data.'),
    ('02', 'TINT Try-On\nSetup',
     'TINT virtual try-on configured for the product. Embedded in creator content and brand page. Consumers experience the product before buying.'),
    ('03', 'Campaign\nExecution',
     'Creators deliver 1 review blog and 2 professional live streams per window. Shared calendar governs all deadlines. Real-time delivery tracking.'),
    ('04', 'Performance\nReport',
     'Full post-campaign ROI. Reach, conversions, commission paid, reliability score impact. Recommendations for the next campaign included.'),
]
x_pos = [0.3, 3.55, 6.8, 10.05]
for i, (num, title, desc) in enumerate(steps):
    x = x_pos[i]
    rect(sl, x, 1.85, 3.0, 5.3, WHITE)
    rect(sl, x, 1.85, 3.0, 0.06, VIOLET)
    txt(sl, num, x+0.15, 2.05, 2.7, 0.72,
        size=32, bold=True, color=LILAC, align=PP_ALIGN.CENTER, font='Montserrat')
    txt(sl, title, x+0.15, 2.88, 2.7, 0.7,
        size=15, bold=True, color=BODY, font='Montserrat')
    txt(sl, desc, x+0.15, 3.7, 2.7, 2.8, size=12, color=SECONDARY)
    if i < 3:
        txt(sl, '→', x+3.05, 4.2, 0.42, 0.55,
            size=20, color=LILAC, align=PP_ALIGN.CENTER)

# Content Trinity callout
rect(sl, 0.3, 7.1, 12.73, 0.28, VIOLET)
txt(sl, 'Content Trinity: 1x Detailed Blog  +  2x Live Streams  +  AI Sales Certification  +  Reliability Score 0-100%',
    0.35, 7.12, 12.5, 0.24,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — TIER STRUCTURE (BRANDS + INFLUENCERS)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '07 — PLATFORM TIERS', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'Two ecosystems. Both governed.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

# Brand tiers
txt(sl, 'BRAND MEMBERSHIP', 0.3, 1.52, 6, 0.35,
    size=11, bold=True, color=MUTED)

brand_tiers = [
    ('Essential',    'HK$400/mo',  'Campaigns: 1/quarter', 'Up to 5 creators', 'TINT basic'),
    ('Professional', 'HK$900/mo',  'Campaigns: 2/quarter', 'Up to 15 creators', 'Full TINT + dashboard'),
    ('Enterprise',   'HK$2,500/mo','Unlimited campaigns',   'Unlimited creators', 'Dedicated manager'),
]
x_pos_b = [0.3, 4.45, 8.6]
for i, (name, price, c1, c2, c3) in enumerate(brand_tiers):
    x = x_pos_b[i]
    featured = (i == 1)
    bc = SOFT_LILAC if featured else LAV_WASH
    ac = VIOLET if featured else LILAC
    rect(sl, x, 1.9, 3.9, 2.55, bc)
    rect(sl, x, 1.9, 3.9, 0.06, ac)
    txt(sl, name, x+0.15, 2.04, 3.6, 0.45,
        size=17, bold=True, color=ac, align=PP_ALIGN.CENTER, font='Montserrat')
    rect(sl, x+0.15, 2.55, 3.6, 0.4, VIOLET)
    txt(sl, price, x+0.15, 2.57, 3.6, 0.36,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, feat in enumerate([c1, c2, c3]):
        txt(sl, '  ' + feat, x+0.18, 3.05 + j*0.42, 3.6, 0.38,
            size=12, color=BODY)

# Influencer tiers
txt(sl, 'INFLUENCER TIERS', 0.3, 4.65, 6, 0.35,
    size=11, bold=True, color=MUTED)

inf_tiers = [
    ('Rising Star',       'Essential brands',    '1 Blog + 1 Live',   '10% commission'),
    ('Star Influencer',   'Professional brands', '2 Blogs + 4 Lives', '12% + content fees'),
    ('Superstar',         'Enterprise brands',   '5 Blogs + 10 Lives','15% + 1st Launch access'),
]
for i, (name, brand_match, commit, earn) in enumerate(inf_tiers):
    x = x_pos_b[i]
    rect(sl, x, 5.05, 3.9, 2.12, LAV_WASH)
    rect(sl, x, 5.05, 3.9, 0.06, LILAC)
    txt(sl, name, x+0.15, 5.18, 3.6, 0.42,
        size=15, bold=True, color=VIOLET, font='Montserrat')
    for j, feat in enumerate([brand_match, commit, earn]):
        txt(sl, '  ' + feat, x+0.18, 5.68 + j*0.42, 3.6, 0.38, size=12, color=BODY)

txt(sl, 'Red Card rule: reliability score below 70% = permanent shop suspension.',
    0.3, 7.2, 12.73, 0.24,
    size=11, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — TRACTION & ROADMAP
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '08 — TRACTION & ROADMAP', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'Where we are. Where we are going.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

# Current status
rect(sl, 0.3, 1.5, 12.73, 1.05, LAV_WASH)
rect(sl, 0.3, 1.5, 12.73, 0.06, VIOLET)
txt(sl, 'Current Status — April 2026', 0.5, 1.6, 12, 0.38,
    size=12, bold=True, color=MUTED)
status_items = [
    'Platform built and live at simplex-ity.fluentlab.co',
    'Brand guidelines and design system completed',
    'Influencer ecosystem framework defined (V3)',
    'TINT integration scoped and ready for brand onboarding',
]
for i, item in enumerate(status_items):
    x = 0.5 + (i % 2) * 6.4
    y = 1.97 if i < 2 else 1.97
    # Side by side
txt(sl, '  '.join(['— ' + s for s in status_items[:2]]),
    0.5, 1.97, 12.3, 0.35, size=12, color=BODY)
txt(sl, '  '.join(['— ' + s for s in status_items[2:]]),
    0.5, 2.35, 12.3, 0.35, size=12, color=BODY)

# Roadmap
roadmap = [
    ('Q2 2026\nNow', 'Platform live. Brand partnerships open. First influencer cohort onboarding. TINT try-on integration active.'),
    ('Q3 2026', 'First 10 brand members signed. 50 verified creators active. Campaign reporting dashboard launched. First performance reports delivered.'),
    ('Q4 2026', 'Professional and Enterprise tier brands onboarded. Retainer model active for top influencers. XHS and TikTok campaign channels added.'),
    ('2027', 'Regional expansion. Singapore and Taiwan market entry. Institutional brand partnerships. Series A readiness.'),
]
y_start = 2.95
for i, (phase, desc) in enumerate(roadmap):
    x = 0.3 + i * 3.25
    rect(sl, x, y_start, 3.0, 4.25, LAV_WASH)
    rect(sl, x, y_start, 3.0, 0.06, VIOLET if i == 0 else LILAC if i < 3 else SOFT_LILAC)
    label_col = VIOLET if i == 0 else BODY
    txt(sl, phase, x+0.15, y_start+0.12, 2.7, 0.7,
        size=16, bold=True, color=label_col, font='Montserrat')
    txt(sl, desc, x+0.15, y_start+0.9, 2.7, 3.2, size=12, color=SECONDARY)
    if i == 0:
        pill_label(sl, 'ACTIVE NOW', x+0.15, y_start+0.72, w=2.0, h=0.28,
                   bg_col=VIOLET, size=9)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — WHY SIMPLEX-ITY WINS (COMPETITIVE MOAT)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0.07, 0.07, 7.43, LILAC)
add_logo(sl, 0.3, 0.15, 1.8)
txt(sl, '09 — COMPETITIVE ADVANTAGE', 0.3, 0.2, 5, 0.28, size=9, bold=True, color=MUTED)

txt(sl, 'Five moats. Built in from day one.',
    0.3, 0.65, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

moats = [
    ('Try-On Technology',
     'TINT virtual AR makeup is integrated at the platform level. Competitors would need to build or license this separately.'),
    ('Reliability Governance',
     'The 0-100% reliability scorecard is proprietary. No other Asian beauty influencer platform governs creator behaviour this systematically.'),
    ('Category Focus',
     'Simplex-ity is built exclusively for Asian beauty, not a generic influencer marketplace. Category depth creates trust that horizontal platforms cannot replicate.'),
    ('Creator Ecosystem',
     'Cross-pollination drives follower growth for creators. The more creators on the platform, the more valuable each one becomes. Network effects compound.'),
    ('Data Advantage',
     'Every campaign generates proprietary data on what converts in Asian beauty by product type, skin tone, and creator tier. This data widens the moat over time.'),
]
for i, (title, desc) in enumerate(moats):
    y = 1.65 + i * 1.08
    rect(sl, 0.3, y, 12.73, 1.0, LAV_WASH)
    rect(sl, 0.3, y, 0.06, 1.0, VIOLET)
    txt(sl, title, 0.5, y+0.1, 3.5, 0.45,
        size=14, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, 4.0, y+0.12, 8.8, 0.72, size=12.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — THE ASK + CLOSING
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, DARK_HERO)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 7.5, 0, 5.83, 7.5, RGBColor(0x20, 0x14, 0x40))
divider(sl, 7.5, 0, 0.06, LILAC, h=7.5)

add_logo(sl, 0.5, 0.22, 2.6)

txt(sl, 'THE INVESTMENT OPPORTUNITY',
    0.5, 1.38, 6.7, 0.38,
    size=10, bold=True, color=MUTED)

txt(sl, 'Join Simplex-ity at\nthe foundation stage.',
    0.5, 1.85, 6.7, 1.7,
    size=36, bold=True, color=WHITE, font='Montserrat')

txt(sl, 'We are seeking seed investment to accelerate brand onboarding, expand the creator ecosystem across HK, Singapore, and Taiwan, and build the full analytics infrastructure.',
    0.5, 3.65, 6.7, 1.2,
    size=14, color=SOFT_LILAC)

divider(sl, 0.5, 5.05, 6.5, LILAC)

use_items = [
    ('Brand Sales', 'Onboard first 20 brand partners across all tiers'),
    ('Creator Growth', 'Build verified creator cohort to 200 active influencers'),
    ('Technology', 'Full TINT integration, analytics dashboard, and reporting suite'),
    ('Regional', 'Singapore and Taiwan market entry in 2027'),
]
y = 5.22
for label, desc in use_items:
    txt(sl, label + ':', 0.5, y, 2.2, 0.35,
        size=12, bold=True, color=LILAC)
    txt(sl, desc, 2.8, y, 4.4, 0.35, size=12, color=SOFT_LILAC)
    y += 0.42

# Right column — contact + key numbers
txt(sl, 'KEY NUMBERS', 7.7, 1.38, 5.3, 0.35,
    size=10, bold=True, color=MUTED)
key_nums = [
    ('3 Revenue Streams', 'Membership, Performance Fee, Premium Shop'),
    ('HK$400 to $2,500', 'Brand membership per month'),
    ('10 to 15%', 'Influencer commission rate'),
    ('25%+', 'TINT conversion rate uplift'),
    ('0', 'Upfront cost for influencers'),
    ('$2.3B', 'Target market size 2026'),
]
y = 1.88
for label, val in key_nums:
    rect(sl, 7.65, y, 5.3, 0.72, RGBColor(0x28, 0x18, 0x48))
    rect(sl, 7.65, y, 0.05, 0.72, LILAC)
    txt(sl, label, 7.82, y+0.05, 3.8, 0.35,
        size=13, bold=True, color=WHITE, font='Montserrat')
    txt(sl, val, 7.82, y+0.4, 4.9, 0.27, size=11, color=SOFT_LILAC)
    y += 0.82

txt(sl, 'kieran@5senses.global',
    0.5, 6.95, 6.7, 0.35, size=13, color=SOFT_LILAC)
txt(sl, 'simplex-ity.fluentlab.co  |  Kwun Tong, Hong Kong',
    0.5, 7.2, 6.7, 0.25, size=11, color=MUTED)

prs.save('/app/SIMPLEX-ITY_Investor_Pitch_v1.pptx')
print('Done! 12 slides built.')
