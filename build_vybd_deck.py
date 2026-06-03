from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
VIOLET = RGBColor(0x5e, 0x50, 0xfb)       # Accent Violet
LILAC  = RGBColor(0x8c, 0x82, 0xfc)       # Primary Lilac
LAVENDER = RGBColor(0xe8, 0xe6, 0xfe)     # Background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1a, 0x1a, 0x1f)
GREY   = RGBColor(0x60, 0x60, 0x70)
LIGHT_VIOLET = RGBColor(0xba, 0xb4, 0xfd)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=LAVENDER):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def txbox(slide, text, l, t, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Montserrat", wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb

def heading_box(slide, text, l, t, w, h, size=28, color=WHITE, font="Exo 2"):
    return txbox(slide, text, l, t, w, h, size=size, bold=True, color=color, align=PP_ALIGN.LEFT, font=font)

def pill(slide, text, l, t, w, h, bg_color=VIOLET, txt_color=WHITE, size=13):
    r = rect(slide, l, t, w, h, bg_color)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = False
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = txt_color
    run.font.name = "Montserrat"
    return r

# ─────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────
s1 = add_slide()
bg(s1, VIOLET)
# Left accent strip
rect(s1, 0, 0, 0.08, 7.5, LILAC)
# Big title
txbox(s1, "SIMPLEX-ITY × VYBD", 0.5, 1.5, 9, 1.2, size=52, bold=True, color=WHITE, font="Exo 2")
txbox(s1, "Strategic Partnership Research Deck", 0.5, 2.75, 9, 0.6, size=22, bold=False, color=LAVENDER, font="Montserrat")
txbox(s1, "Why Vybd is the US Operations Backbone for SIM", 0.5, 3.35, 9, 0.5, size=17, bold=False, color=LIGHT_VIOLET, font="Montserrat")
# Bottom tag
rect(s1, 0, 6.7, 13.33, 0.8, RGBColor(0x4a, 0x3f, 0xd4))
txbox(s1, "Prepared by Simpee  |  June 2026  |  CONFIDENTIAL", 0.5, 6.75, 12, 0.5, size=13, color=LAVENDER, font="Montserrat")
# Right logo placeholder text
txbox(s1, "SIM", 11.0, 1.5, 2, 1.0, size=48, bold=True, color=LAVENDER, align=PP_ALIGN.CENTER, font="Exo 2")
txbox(s1, "× VYBD", 10.8, 2.5, 2.3, 0.5, size=18, color=LIGHT_VIOLET, align=PP_ALIGN.CENTER, font="Montserrat")

# ─────────────────────────────────────────────
# SLIDE 2 — The Problem SIM Faces Without Vybd
# ─────────────────────────────────────────────
s2 = add_slide()
bg(s2)
rect(s2, 0, 0, 13.33, 1.2, VIOLET)
txbox(s2, "THE PROBLEM  —  Without Vybd, SIM Hits a Wall", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

problems = [
    ("01", "No US Logistics",       "How do Korean skincare brands get products to US influencers? SIM has no answer without a fulfilment partner."),
    ("02", "FDA & Customs Gap",     "Every cosmetics product entering the US needs FDA registration & customs classification. One mistake = shipment blocked."),
    ("03", "Scattered Vendors",     "Without Vybd, SIM would need to coordinate 3PLs, customs brokers, agencies & software separately — impossible at launch."),
    ("04", "Brand Trust Problem",   "Asian brands won't commit to SIM if SIM can't explain how their products physically reach US influencers."),
]

cols = [0.3, 3.5, 6.7, 9.9]
for i, (num, title, body) in enumerate(problems):
    x = cols[i]
    rect(s2, x, 1.4, 3.0, 5.6, WHITE)
    rect(s2, x, 1.4, 3.0, 0.65, LILAC)
    txbox(s2, num, x+0.1, 1.42, 0.6, 0.6, size=28, bold=True, color=WHITE, font="Exo 2")
    txbox(s2, title, x+0.65, 1.45, 2.2, 0.6, size=15, bold=True, color=WHITE, font="Exo 2")
    txbox(s2, body, x+0.15, 2.2, 2.7, 4.5, size=13, color=DARK, font="Montserrat")

# ─────────────────────────────────────────────
# SLIDE 3 — What Vybd Is
# ─────────────────────────────────────────────
s3 = add_slide()
bg(s3)
rect(s3, 0, 0, 13.33, 1.2, DARK)
txbox(s3, "WHAT IS VYBD?", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

txbox(s3, '"You bring the vibe. We bring the operations team."', 0.4, 1.35, 12, 0.7, size=20, bold=False, color=VIOLET, align=PP_ALIGN.CENTER, font="Montserrat")
txbox(s3, "Vybd is an AI-powered US market operations OS built for international brands entering the US market.\nTarget: Brands with $500K–$5M revenue in home market. Promise: 14-day US market launch.", 0.4, 2.1, 12.5, 0.9, size=14, color=DARK, align=PP_ALIGN.CENTER, font="Montserrat")

# 7 agent boxes
agents = [
    ("🛃", "Compliance",       "Customs docs, FDA"),
    ("📊", "Market Intel",     "Competitors, pricing"),
    ("🛒", "E-Commerce",       "Listings, inventory"),
    ("📣", "Marketing",        "Campaigns, content"),
    ("🚚", "Logistics",        "Carriers, tracking"),
    ("🏭", "Warehousing",      "Stock, replenishment"),
    ("⚙️", "Operations",       "Coordination, exceptions"),
]

positions = [0.3, 2.2, 4.1, 6.0, 7.9, 9.8, 11.7]
for i, (icon, name, desc) in enumerate(agents):
    x = positions[i]
    rect(s3, x, 3.1, 1.7, 3.8, WHITE)
    rect(s3, x, 3.1, 1.7, 0.7, VIOLET)
    txbox(s3, icon, x+0.55, 3.12, 0.6, 0.65, size=20, color=WHITE, font="Segoe UI Emoji")
    txbox(s3, name, x+0.05, 3.85, 1.6, 0.5, size=12, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Exo 2")
    txbox(s3, desc, x+0.08, 4.45, 1.55, 2.2, size=11, color=GREY, align=PP_ALIGN.CENTER, font="Montserrat")

txbox(s3, "vybd.ai  |  hello@vybd.com  |  Status: Early Access Pilot (Limited Spots)", 0.4, 7.0, 12.5, 0.4, size=11, color=GREY, align=PP_ALIGN.CENTER, font="Montserrat")

# ─────────────────────────────────────────────
# SLIDE 4 — How SIM + Vybd Work Together
# ─────────────────────────────────────────────
s4 = add_slide()
bg(s4)
rect(s4, 0, 0, 13.33, 1.2, LILAC)
txbox(s4, "HOW SIM + VYBD WORK TOGETHER", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

# SIM column
rect(s4, 0.3, 1.4, 5.5, 5.7, WHITE)
rect(s4, 0.3, 1.4, 5.5, 0.65, VIOLET)
txbox(s4, "SIM HANDLES", 0.5, 1.45, 5.0, 0.55, size=17, bold=True, color=WHITE, font="Exo 2")
sim_items = [
    "Brand Recruitment — onboard Korean skincare brands",
    "Influencer Matching — connect brands with US influencers",
    "Campaign Strategy — set objectives & content briefs",
    "Platform Data — capture performance metrics",
    "NEST VC Reporting — present validated trial data",
]
for j, item in enumerate(sim_items):
    txbox(s4, f"✓  {item}", 0.5, 2.2 + j*0.82, 5.1, 0.75, size=13, color=DARK, font="Montserrat")

# Arrow
txbox(s4, "⟷", 5.95, 3.8, 1.3, 0.9, size=36, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Segoe UI Emoji")
txbox(s4, "FULL\nPARTNER", 5.95, 4.65, 1.3, 0.7, size=11, bold=True, color=GREY, align=PP_ALIGN.CENTER, font="Montserrat")

# Vybd column
rect(s4, 7.5, 1.4, 5.5, 5.7, WHITE)
rect(s4, 7.5, 1.4, 5.5, 0.65, DARK)
txbox(s4, "VYBD HANDLES", 7.7, 1.45, 5.0, 0.55, size=17, bold=True, color=WHITE, font="Exo 2")
vybd_items = [
    "US Customs & FDA compliance per brand",
    "US Warehousing for product samples",
    "Shipping samples to US influencers",
    "US Market Intelligence per brand",
    "Operations coordination & exceptions",
]
for j, item in enumerate(vybd_items):
    txbox(s4, f"✓  {item}", 7.7, 2.2 + j*0.82, 5.1, 0.75, size=13, color=DARK, font="Montserrat")

txbox(s4, "Result: SIM becomes a TRUE end-to-end US market entry platform — not just a matchmaker.", 0.4, 7.05, 12.5, 0.4, size=13, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Montserrat")

# ─────────────────────────────────────────────
# SLIDE 5 — Key Benefits to SIM
# ─────────────────────────────────────────────
s5 = add_slide()
bg(s5)
rect(s5, 0, 0, 13.33, 1.2, VIOLET)
txbox(s5, "KEY BENEFITS TO SIM", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

benefits = [
    ("Zero Headcount\nfor US Ops", "No need to hire logistics, compliance or ops staff in the US. Vybd's AI agents do it all."),
    ("14-Day Brand\nUS Launch", "Each new brand SIM onboards can be live in the US within 14 days — powerful pitch to NEST VC."),
    ("Full NEST VC\nStory", "SIM + Vybd = complete platform with validated ops. Not just a concept — a working system."),
    ("Scales with\nZero SIM Overhead", "Add 10 or 100 brands — Vybd scales. SIM's cost per brand drops as volume grows."),
    ("Brand Trust\nInstantly", "Korean brands ask 'how do products reach influencers?' Vybd is the answer. Removes biggest objection."),
    ("Co-Pilot\nOpportunity", "Both in early stage. SIM can negotiate a strategic co-pilot deal — not just a customer fee."),
]

cols2 = [0.3, 4.55, 8.8]
rows2 = [1.4, 4.2]
idx = 0
for row in rows2:
    for col in cols2:
        if idx >= len(benefits): break
        title, body = benefits[idx]
        rect(s5, col, row, 3.9, 2.5, WHITE)
        rect(s5, col, row, 3.9, 0.08, VIOLET)
        txbox(s5, title, col+0.15, row+0.15, 3.6, 0.85, size=15, bold=True, color=VIOLET, font="Exo 2")
        txbox(s5, body, col+0.15, row+1.05, 3.6, 1.3, size=12, color=DARK, font="Montserrat")
        idx += 1

# ─────────────────────────────────────────────
# SLIDE 6 — Estimated Costs
# ─────────────────────────────────────────────
s6 = add_slide()
bg(s6)
rect(s6, 0, 0, 13.33, 1.2, DARK)
txbox(s6, "ESTIMATED COSTS  —  3-Month Trial (5 Brands)", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

# Note: Vybd pricing not yet public — estimates based on comparable platforms
txbox(s6, "Note: Vybd pricing is not yet public (pilot stage). Estimates based on comparable US fulfilment + ops platforms.", 0.4, 1.3, 12.5, 0.45, size=12, bold=False, color=GREY, align=PP_ALIGN.CENTER, font="Montserrat")

# Table header
rect(s6, 0.3, 1.85, 12.7, 0.55, VIOLET)
txbox(s6, "COST ITEM", 0.4, 1.88, 4.5, 0.5, size=13, bold=True, color=WHITE, font="Exo 2")
txbox(s6, "UNIT", 4.9, 1.88, 2.5, 0.5, size=13, bold=True, color=WHITE, font="Exo 2")
txbox(s6, "EST. RANGE (USD)", 7.4, 1.88, 3.5, 0.5, size=13, bold=True, color=WHITE, font="Exo 2")
txbox(s6, "TYPE", 10.9, 1.88, 2.0, 0.5, size=13, bold=True, color=WHITE, font="Exo 2")

rows_data = [
    ("Vybd Platform / Ops Fee",     "Per brand / month",  "$800 – $2,000",    "Monthly"),
    ("US Customs & FDA Setup",      "Per brand",          "$500 – $1,500",    "One-time"),
    ("US Warehousing",              "Per brand / month",  "$200 – $500",      "Monthly"),
    ("Influencer Sample Shipping",  "Per package (US)",   "$8 – $25",         "Per shipment"),
    ("Market Intelligence Reports", "Per brand / month",  "$100 – $300",      "Monthly"),
]

for i, (item, unit, est, typ) in enumerate(rows_data):
    row_y = 2.55 + i * 0.72
    fill = LAVENDER if i % 2 == 0 else WHITE
    rect(s6, 0.3, row_y, 12.7, 0.68, fill)
    txbox(s6, item, 0.4, row_y+0.08, 4.5, 0.55, size=13, color=DARK, font="Montserrat")
    txbox(s6, unit, 4.9, row_y+0.08, 2.5, 0.55, size=13, color=GREY, font="Montserrat")
    txbox(s6, est,  7.4, row_y+0.08, 3.5, 0.55, size=13, bold=True, color=VIOLET, font="Montserrat")
    txbox(s6, typ,  10.9, row_y+0.08, 2.0, 0.55, size=12, color=DARK, font="Montserrat")

# Summary box
rect(s6, 0.3, 6.3, 12.7, 0.9, VIOLET)
txbox(s6, "5 Brands — Monthly Burn Estimate:", 0.5, 6.35, 5.5, 0.7, size=14, bold=True, color=WHITE, font="Exo 2")
txbox(s6, "$5,500 – $14,500 / month", 6.0, 6.35, 7.0, 0.7, size=20, bold=True, color=LAVENDER, font="Exo 2")

# ─────────────────────────────────────────────
# SLIDE 7 — Revenue Potential
# ─────────────────────────────────────────────
s7 = add_slide()
bg(s7)
rect(s7, 0, 0, 13.33, 1.2, LILAC)
txbox(s7, "PREDICTED VALUE  —  Revenue Potential Post-Trial", 0.4, 0.2, 12, 0.8, size=26, bold=True, color=WHITE, font="Exo 2")

# 3 scenario boxes
scenarios = [
    ("CONSERVATIVE", "10 Brands", "$2,000/mo platform fee\n+ 15% campaign commission", "$20,000 – $25,000 MRR"),
    ("MODERATE",     "25 Brands", "$2,000/mo platform fee\n+ 15% campaign commission", "$50,000 – $65,000 MRR"),
    ("OPTIMISTIC",   "50 Brands", "$2,000/mo platform fee\n+ 15% campaign commission", "$100,000 – $130,000 MRR"),
]

s_cols = [0.5, 4.7, 8.9]
s_colors = [LILAC, VIOLET, DARK]
for i, (label, brands, model, mrr) in enumerate(scenarios):
    x = s_cols[i]
    rect(s7, x, 1.4, 3.9, 5.3, WHITE)
    rect(s7, x, 1.4, 3.9, 0.7, s_colors[i])
    txbox(s7, label, x+0.15, 1.45, 3.6, 0.6, size=16, bold=True, color=WHITE, font="Exo 2")
    txbox(s7, brands, x+0.15, 2.3, 3.6, 0.5, size=22, bold=True, color=s_colors[i], font="Exo 2")
    txbox(s7, model, x+0.15, 2.95, 3.6, 0.9, size=12, color=GREY, font="Montserrat")
    rect(s7, x+0.15, 4.05, 3.6, 1.2, LAVENDER)
    txbox(s7, mrr, x+0.25, 4.15, 3.4, 1.0, size=17, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Exo 2")

txbox(s7, "With Vybd as ops backbone, SIM can scale from 5 trial brands to 50 without adding headcount.", 0.4, 6.85, 12.5, 0.5, size=13, bold=False, color=DARK, align=PP_ALIGN.CENTER, font="Montserrat")

txbox(s7, "These figures represent SIM platform revenue only. Vybd fees reduce net margin but are largely passed to brands.", 0.4, 7.1, 12.5, 0.35, size=11, color=GREY, align=PP_ALIGN.CENTER, font="Montserrat")

# ─────────────────────────────────────────────
# SLIDE 8 — Why Now / Next Steps
# ─────────────────────────────────────────────
s8 = add_slide()
bg(s8, VIOLET)
rect(s8, 0, 0, 0.08, 7.5, LAVENDER)
txbox(s8, "WHY NOW  &  NEXT STEPS", 0.4, 0.3, 12, 0.9, size=32, bold=True, color=WHITE, font="Exo 2")

why_now = [
    "Both SIM and Vybd are in early/pilot stage — leverage available",
    "Vybd has limited pilot spots — first mover advantage",
    "NEST VC needs a complete platform story — Vybd fills the ops gap",
    "SIM's 3-month trial plan needs US fulfilment NOW — not later",
]

txbox(s8, "WHY NOW", 0.5, 1.35, 5.8, 0.5, size=16, bold=True, color=LAVENDER, font="Exo 2")
for j, pt in enumerate(why_now):
    rect(s8, 0.5, 1.95 + j*0.85, 5.8, 0.72, RGBColor(0x4a, 0x3f, 0xd4))
    txbox(s8, f"→  {pt}", 0.65, 2.0 + j*0.85, 5.5, 0.65, size=12, color=WHITE, font="Montserrat")

steps = [
    ("Week 1", "Email Vybd at hello@vybd.com — frame as strategic co-pilot, not just customer"),
    ("Week 2", "Schedule intro call — present SIM's 3-month trial plan + brand pipeline"),
    ("Week 3", "Negotiate pilot pricing — aim for revenue-share model vs fixed monthly fee"),
    ("Week 4", "Sign MOU — Vybd confirmed as SIM's official US Ops Partner"),
]

txbox(s8, "NEXT STEPS", 7.3, 1.35, 5.7, 0.5, size=16, bold=True, color=LAVENDER, font="Exo 2")
for j, (week, action) in enumerate(steps):
    rect(s8, 7.3, 1.95 + j*0.85, 1.1, 0.72, LAVENDER)
    txbox(s8, week, 7.32, 2.0 + j*0.85, 1.05, 0.65, size=12, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Exo 2")
    rect(s8, 8.45, 1.95 + j*0.85, 4.5, 0.72, RGBColor(0x4a, 0x3f, 0xd4))
    txbox(s8, action, 8.55, 2.0 + j*0.85, 4.3, 0.65, size=12, color=WHITE, font="Montserrat")

txbox(s8, "Prepared by Simpee  ×  SIMPLEX-ITY  |  June 2026  |  CONFIDENTIAL  |  vybd.ai", 0.4, 7.05, 12.5, 0.4, size=11, color=LAVENDER, align=PP_ALIGN.CENTER, font="Montserrat")

# ─────────────────────────────────────────────
prs.save("/app/SIM_x_Vybd_Partnership_Deck.pptx")
print("DONE")
