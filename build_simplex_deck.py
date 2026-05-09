from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1a, 0x05, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

# ─── SLIDE 1: COVER ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
# Purple accent bar left
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
# Decorative circle
shape = slide.shapes.add_shape(9, Inches(9), Inches(-1), Inches(5), Inches(5))
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE
shape.line.fill.background()

add_text(slide, "SIMPLEX-ITY", 1, 1.5, 9, 1.2, size=54, bold=True, color=WHITE)
add_text(slide, "Beauty-Tech Marketing Platform", 1, 2.9, 9, 0.7, size=24, color=LIGHT_PURPLE, italic=True)
add_rect(slide, 1, 3.8, 4, 0.06, PURPLE)
add_text(slide, "Brand Partnership Deck  |  2026", 1, 4.0, 8, 0.6, size=16, color=GRAY)
add_text(slide, "Connecting Brands × Influencers × AI Try-On Technology", 1, 5.2, 11, 0.7, size=18, color=WHITE)
add_text(slide, "kieran@5senses.global", 1, 6.3, 6, 0.5, size=13, color=GRAY)

# ─── SLIDE 2: THE PROBLEM ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "THE PROBLEM", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "Beauty brands are wasting budget on the wrong influencers", 1, 1.0, 11, 1.0, size=32, bold=True, color=WHITE)

problems = [
    ("❌", "No way to verify if an influencer's audience actually buys beauty products"),
    ("❌", "Virtual try-on technology is expensive & disconnected from influencer campaigns"),
    ("❌", "HK beauty brands have no centralised platform to manage influencer ROI"),
    ("❌", "Influencer reliability is tracked manually — campaigns fall apart mid-execution"),
]
y = 2.4
for icon, text in problems:
    add_rect(slide, 1, y, 11, 0.55, RGBColor(0x2D, 0x10, 0x4A))
    add_text(slide, icon + "  " + text, 1.2, y + 0.05, 10.5, 0.5, size=16, color=WHITE)
    y += 0.72

# ─── SLIDE 3: THE SOLUTION ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "THE SOLUTION", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "SIMPLEX-ITY — The All-in-One Beauty Marketing Platform", 1, 1.0, 11, 1.0, size=28, bold=True, color=WHITE)

solutions = [
    ("🤖", "AI-Powered Matching", "Smart algorithm pairs your brand with the most relevant influencers based on audience data and past performance"),
    ("💄", "TINT AI Try-On", "Integrated virtual makeup try-on so consumers can test your products before buying — directly in influencer content"),
    ("📊", "Reliability Scoring", "Every influencer is rated on delivery, engagement, and conversion — so you only work with proven performers"),
    ("🚀", "End-to-End Campaigns", "From brief to delivery, we manage the entire campaign lifecycle — live counts, blog posts, commission tracking"),
]
y = 2.2
for icon, title, desc in solutions:
    add_rect(slide, 1, y, 5.3, 1.0, RGBColor(0x2D, 0x10, 0x4A))
    add_text(slide, icon + " " + title, 1.2, y + 0.05, 5, 0.4, size=15, bold=True, color=PURPLE)
    add_text(slide, desc, 1.2, y + 0.45, 5, 0.5, size=12, color=WHITE)
    y += 1.15
    if y > 3.4:
        y = 2.2
        # shift to right column after 2
        # handled by layout below

# Redo layout as 2x2 grid
for shape in slide.shapes:
    if hasattr(shape, "text") and shape.text not in ["THE SOLUTION", "SIMPLEX-ITY — The All-in-One Beauty Marketing Platform"]:
        sp = shape._element
        sp.getparent().remove(sp)

cols = [(1, 2.3), (7, 2.3), (1, 4.7), (7, 4.7)]
for i, (icon, title, desc) in enumerate(solutions):
    lx, ly = cols[i]
    add_rect(slide, lx, ly, 5.5, 2.0, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, lx, ly, 5.5, 0.45, PURPLE)
    add_text(slide, icon + "  " + title, lx + 0.15, ly + 0.05, 5.2, 0.38, size=15, bold=True, color=WHITE)
    add_text(slide, desc, lx + 0.15, ly + 0.55, 5.1, 1.3, size=13, color=WHITE)

# ─── SLIDE 4: HOW IT WORKS ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "HOW IT WORKS", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "Simple 4-Step Campaign Journey", 1, 1.0, 11, 0.7, size=30, bold=True, color=WHITE)

steps = [
    ("01", "Brand Onboards", "Share your product, budget & target audience. We match you to verified influencers instantly."),
    ("02", "TINT Try-On Setup", "We integrate AI virtual try-on into your campaign so consumers can try before they buy."),
    ("03", "Campaign Execution", "Influencers go LIVE and publish blogs. We track every delivery in real time."),
    ("04", "Performance Report", "Full ROI report — reach, engagement, conversions, commission breakdown."),
]
x_positions = [1, 3.9, 6.8, 9.7]
for i, (num, title, desc) in enumerate(steps):
    x = x_positions[i]
    add_rect(slide, x, 2.0, 2.7, 4.5, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, 2.0, 2.7, 0.7, PURPLE)
    add_text(slide, num, x + 0.1, 2.05, 0.8, 0.6, size=28, bold=True, color=WHITE)
    add_text(slide, title, x + 0.15, 2.85, 2.4, 0.55, size=15, bold=True, color=WHITE)
    add_text(slide, desc, x + 0.15, 3.5, 2.4, 2.8, size=12, color=LIGHT_PURPLE)
    # Arrow between steps
    if i < 3:
        add_text(slide, "→", x + 2.75, 3.8, 0.4, 0.5, size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ─── SLIDE 5: TINT AI TECHNOLOGY ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "OUR EDGE", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "TINT AI Virtual Try-On Technology", 1, 1.0, 11, 0.7, size=30, bold=True, color=WHITE)
add_text(slide, "Powered by Banuba — the world's most realistic makeup try-on engine", 1, 1.85, 11, 0.5, size=16, color=LIGHT_PURPLE, italic=True)

stats = [
    ("200%", "More Conversions"),
    ("30%", "Higher Avg Order Value"),
    ("300%", "Higher Engagement Rate"),
    ("17%", "AI Beauty Market CAGR"),
]
x_positions = [1, 4, 7, 10]
for i, (num, label) in enumerate(stats):
    x = x_positions[i]
    add_rect(slide, x, 2.6, 2.5, 2.0, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, 2.6, 2.5, 0.06, PURPLE)
    add_text(slide, num, x + 0.1, 2.75, 2.3, 0.9, size=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 3.65, 2.3, 0.7, size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "What TINT Does:", 1, 4.9, 11, 0.45, size=16, bold=True, color=WHITE)
tint_points = "✅  Real-time AR lipstick, eyeshadow, foundation try-on   ✅  AI skin tone analysis & product recommendation   ✅  Embeds directly in influencer content & brand pages"
add_text(slide, tint_points, 1, 5.45, 11, 0.9, size=14, color=LIGHT_PURPLE)

# ─── SLIDE 6: MARKET OPPORTUNITY ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "MARKET OPPORTUNITY", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "A Market Ready to Be Captured", 1, 1.0, 11, 0.7, size=30, bold=True, color=WHITE)

market_data = [
    ("$2.3B", "Global AI Beauty\nPersonalization Market\n2026", PURPLE),
    ("$2.03B", "Influencer Marketing\nPlatform Market\nby 2031", RGBColor(0x06, 0x95, 0x6E)),
    ("45.3%", "Instagram's share of\nHK Influencer\nMarketing 2025", GOLD),
    ("68%", "Beauty purchases\ninfluenced by\nsocial discovery", RGBColor(0xEF, 0x44, 0x44)),
]
x_positions = [1, 4, 7, 10]
for i, (num, label, color) in enumerate(market_data):
    x = x_positions[i]
    add_rect(slide, x, 2.0, 2.5, 3.0, RGBColor(0x2D, 0x10, 0x4A))
    add_rect(slide, x, 2.0, 2.5, 0.08, color)
    add_text(slide, num, x + 0.1, 2.2, 2.3, 0.9, size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 3.2, 2.3, 1.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "HK Specific Insight:", 1, 5.3, 11, 0.4, size=15, bold=True, color=WHITE)
add_text(slide, "Instagram dominates (45.3%) but XHS (Xiaohongshu) is the fastest growing platform in HK beauty — SIMPLEX-ITY is positioned across BOTH.", 1, 5.8, 11, 0.8, size=14, color=LIGHT_PURPLE)

# ─── SLIDE 7: WHAT YOU GET ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "BRAND PARTNERSHIP", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "What You Get as a SIMPLEX-ITY Brand Partner", 1, 1.0, 11, 0.7, size=28, bold=True, color=WHITE)

deliverables = [
    ("🎯", "Verified Influencer Matching", "AI-curated list of influencers whose audience matches your brand's target customer"),
    ("💄", "TINT Try-On Integration", "Virtual makeup try-on embedded in influencer posts and your brand page"),
    ("📱", "Multi-Platform Campaigns", "Instagram, XHS, TikTok — coordinated across platforms with unified tracking"),
    ("📊", "Real-Time Dashboard", "Live campaign performance — views, clicks, conversions, commission paid"),
    ("🛡️", "Reliability Guarantee", "All influencers are pre-vetted with reliability scores — no ghost deliveries"),
    ("📋", "Full Campaign Report", "Post-campaign ROI analysis with recommendations for next campaign"),
]
cols = [(1, 2.1), (7, 2.1), (1, 3.45), (7, 3.45), (1, 4.8), (7, 4.8)]
for i, (icon, title, desc) in enumerate(deliverables):
    lx, ly = cols[i]
    add_text(slide, icon + "  " + title, lx, ly, 5.5, 0.45, size=15, bold=True, color=PURPLE)
    add_text(slide, desc, lx + 0.35, ly + 0.4, 5.5, 0.6, size=13, color=WHITE)

# ─── SLIDE 8: MEMBERSHIP TIERS ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
add_text(slide, "PRICING", 1, 0.4, 10, 0.7, size=13, bold=True, color=PURPLE)
add_text(slide, "Brand Membership Tiers", 1, 1.0, 11, 0.7, size=30, bold=True, color=WHITE)

tiers = [
    ("STARTER", "Entry-level brands\n& new launches", ["1 campaign/quarter", "Up to 5 influencers", "Basic TINT try-on", "Performance report"], RGBColor(0x6B, 0x72, 0x80), "Contact for pricing"),
    ("GROWTH", "Scaling brands\n& seasonal pushes", ["2 campaigns/quarter", "Up to 15 influencers", "Full TINT integration", "Real-time dashboard", "Priority matching"], PURPLE, "Contact for pricing"),
    ("PREMIUM", "Established brands\n& multi-channel", ["Unlimited campaigns", "Unlimited influencers", "Custom TINT features", "Dedicated manager", "Monthly strategy call"], GOLD, "Contact for pricing"),
]
x_positions = [0.8, 4.7, 8.6]
for i, (name, target, features, color, price) in enumerate(tiers):
    x = x_positions[i]
    is_featured = i == 1
    add_rect(slide, x, 1.9, 3.8, 5.2, RGBColor(0x2D, 0x10, 0x4A) if not is_featured else RGBColor(0x3D, 0x15, 0x6A))
    add_rect(slide, x, 1.9, 3.8, 0.55, color)
    add_text(slide, name, x + 0.1, 1.95, 3.6, 0.45, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, target, x + 0.1, 2.6, 3.6, 0.7, size=12, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER, italic=True)
    y_feat = 3.45
    for feat in features:
        add_text(slide, "✓  " + feat, x + 0.2, y_feat, 3.4, 0.38, size=12, color=WHITE)
        y_feat += 0.38
    add_rect(slide, x + 0.3, 6.6, 3.2, 0.38, color)
    add_text(slide, price, x + 0.3, 6.62, 3.2, 0.35, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 9: CTA / CONTACT ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK)
add_rect(slide, 0, 0, 0.5, 7.5, PURPLE)
# Large purple accent
shape = slide.shapes.add_shape(9, Inches(8), Inches(3), Inches(8), Inches(8))
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE
shape.line.fill.background()

add_text(slide, "Let's Build Something\nBeautiful Together.", 1, 1.2, 8, 2.0, size=38, bold=True, color=WHITE)
add_text(slide, "Join SIMPLEX-ITY as a Brand Partner", 1, 3.5, 8, 0.6, size=20, color=LIGHT_PURPLE)
add_rect(slide, 1, 4.3, 5, 0.06, PURPLE)

contact_info = [
    "📧  kieran@5senses.global",
    "🏢  Room 1608, 16/F APEC Plaza, 49 Hoi Yuen Road, Kwun Tong, HK",
    "🌐  SIMPLEX-ITY by 5SENSESBEAUTY LIMITED",
]
y = 4.6
for info in contact_info:
    add_text(slide, info, 1, y, 9, 0.45, size=15, color=WHITE)
    y += 0.55

add_text(slide, "BR No: 78459506-001-07-25-A", 1, 6.3, 8, 0.4, size=12, color=GRAY, italic=True)

prs.save('/app/SIMPLEX-ITY_Brand_Pitch_Deck.pptx')
print("Deck saved!")
