from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Official Brand Colors ─────────────────────────────────────────
LILAC       = RGBColor(0x8c, 0x82, 0xfc)   # Primary Lilac
VIOLET      = RGBColor(0x5e, 0x50, 0xfb)   # Accent Violet
SOFT_LILAC  = RGBColor(0xba, 0xb4, 0xfd)   # Soft Lilac
LAV_WASH    = RGBColor(0xe8, 0xe6, 0xfe)   # Lavender Wash
WHITE       = RGBColor(0xff, 0xff, 0xff)   # White Canvas
NEUTRAL     = RGBColor(0xe6, 0xe6, 0xe6)   # Neutral Grey
BODY        = RGBColor(0x1a, 0x1a, 0x1f)   # Body Text
MUTED       = RGBColor(0x99, 0x99, 0xaa)   # Text Muted
SECONDARY   = RGBColor(0x5a, 0x5a, 0x6a)   # Text Secondary

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

LOGO_PATH = '/app/logo.jpg'

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False, font='Montserrat'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    r.font.name = font
    return tb

def logo(slide, l, t, w):
    slide.shapes.add_picture(LOGO_PATH, Inches(l), Inches(t), width=Inches(w))

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER  (White canvas, lilac accents)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)

# Top accent bar (violet)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
# Lavender wash right-side panel
rect(sl, 8.5, 0, 4.83, 7.5, LAV_WASH)
# Soft lilac decorative block
rect(sl, 8.5, 4.5, 4.83, 3.0, SOFT_LILAC)

# Logo top-left
logo(sl, 0.6, 0.35, 3.2)

# Eyebrow
txt(sl, 'BRAND PARTNERSHIP DECK  ·  2026', 0.6, 1.8, 7, 0.4,
    size=10, bold=True, color=MUTED)

# Hero headline — Exo 2 style
txt(sl, 'Asian beauty,', 0.6, 2.25, 7.5, 1.0,
    size=48, bold=True, color=BODY, font='Montserrat')
txt(sl, 'curated with intent.', 0.6, 3.1, 7.5, 1.0,
    size=48, bold=True, color=LILAC, font='Montserrat')

# Subline
txt(sl, 'Where creators guide discovery — and brands connect\nwith audiences who are ready to buy.',
    0.6, 4.3, 7.2, 0.9, size=16, color=SECONDARY)

# Tagline pill
rect(sl, 0.6, 5.5, 5.5, 0.48, VIOLET)
txt(sl, 'ONE STOP ASIAN BEAUTY  :  SIMPLIFY TO AMPLIFY',
    0.65, 5.54, 5.4, 0.4, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, 'kieran@5senses.global', 0.6, 6.6, 5, 0.4, size=12, color=MUTED)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0, 0.07, 7.5, LILAC)

logo(sl, 0.3, 0.18, 1.8)
txt(sl, '01 — THE PROBLEM', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'The beauty brand challenge.', 0.3, 0.7, 12, 0.75,
    size=34, bold=True, color=BODY, font='Montserrat')
txt(sl, 'Too many influencers. Too little certainty. No way to see real results.',
    0.3, 1.55, 10, 0.45, size=16, color=SECONDARY, italic=True)

problems = [
    ('No way to verify if an influencer\'s audience actually buys beauty products.'),
    ('Virtual try-on technology is expensive and disconnected from campaigns.'),
    ('HK beauty brands have no centralised platform to manage influencer ROI.'),
    ('Influencer reliability tracked manually — campaigns fail mid-execution.'),
]
y = 2.3
for p in problems:
    rect(sl, 0.3, y, 12.7, 0.78, LAV_WASH)
    rect(sl, 0.3, y, 0.07, 0.78, VIOLET)
    txt(sl, p, 0.55, y + 0.18, 12.0, 0.5, size=15, color=BODY)
    y += 0.97

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0, 0.07, 7.5, LILAC)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '02 — THE SOLUTION', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'Simplex-ity — beauty marketing, simplified.', 0.3, 0.7, 12, 0.75,
    size=32, bold=True, color=BODY, font='Montserrat')
txt(sl, 'One platform. Verified creators. Real results.',
    0.3, 1.5, 10, 0.4, size=16, color=SECONDARY, italic=True)

solutions = [
    ('Verified Creator Matching',
     'Curated influencers matched to your brand by audience data and delivery record — not follower count.'),
    ('TINT Virtual Try-On',
     'Integrated makeup try-on so consumers experience your products inside creator content, before they buy.'),
    ('Reliability Scoring',
     'Every creator is rated on delivery, engagement, and conversion. You only work with proven performers.'),
    ('End-to-End Campaign Management',
     'From brief to full performance report — live counts, blog posts, commission tracking in one place.'),
]
cols = [(0.3, 2.2), (6.8, 2.2), (0.3, 4.65), (6.8, 4.65)]
for i, (title, desc) in enumerate(solutions):
    x, y = cols[i]
    rect(sl, x, y, 6.2, 2.15, LAV_WASH)
    rect(sl, x, y, 6.2, 0.07, VIOLET)
    txt(sl, title, x + 0.2, y + 0.18, 5.8, 0.45, size=15, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, x + 0.2, y + 0.72, 5.8, 1.2, size=13, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '03 — HOW IT WORKS', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'Four steps from brief to results.', 0.3, 0.7, 12, 0.65,
    size=32, bold=True, color=BODY, font='Montserrat')

steps = [
    ('01', 'Brief & Match',
     'Share your product, budget and audience. We match you to creators selected for real relevance.'),
    ('02', 'Try-On Setup',
     'TINT virtual makeup try-on embedded in creator content — consumers experience before they buy.'),
    ('03', 'Campaign Live',
     'Creators go live and publish. Every delivery tracked in real time — live counts, posts, reach.'),
    ('04', 'Performance Report',
     'Full ROI report — reach, engagement, conversions, commission breakdown, next steps.'),
]
x_pos = [0.3, 3.55, 6.8, 10.05]
for i, (num, title, desc) in enumerate(steps):
    x = x_pos[i]
    rect(sl, x, 1.85, 3.0, 5.3, WHITE)
    rect(sl, x, 1.85, 3.0, 0.07, VIOLET)
    txt(sl, num, x + 0.2, 2.05, 1, 0.65, size=30, bold=True, color=LILAC, font='Montserrat')
    txt(sl, title, x + 0.2, 2.8, 2.6, 0.5, size=15, bold=True, color=BODY, font='Montserrat')
    txt(sl, desc, x + 0.2, 3.45, 2.6, 2.5, size=12, color=SECONDARY)
    if i < 3:
        txt(sl, '→', x + 3.05, 3.8, 0.4, 0.5, size=18, color=LILAC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — TINT TECHNOLOGY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0, 0.07, 7.5, LILAC)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '04 — WHAT SETS US APART', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'TINT Virtual Try-On.', 0.3, 0.7, 8, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')
txt(sl, 'The most realistic makeup try-on available — embedded directly in creator content.',
    0.3, 1.45, 10, 0.45, size=15, color=SECONDARY, italic=True)

stats = [
    ('200%', 'More conversions\nvs standard images'),
    ('30%',  'Higher average\norder value'),
    ('300%', 'Higher engagement\nrate'),
    ('17%',  'AI beauty market\nCAGR 2025–2033'),
]
for i, (val, label) in enumerate(stats):
    x = 0.3 + i * 3.25
    rect(sl, x, 2.15, 3.0, 2.2, LAV_WASH)
    rect(sl, x, 2.15, 3.0, 0.07, VIOLET)
    txt(sl, val, x + 0.15, 2.3, 2.7, 0.9,
        size=36, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font='Montserrat')
    txt(sl, label, x + 0.15, 3.25, 2.7, 0.9,
        size=13, color=BODY, align=PP_ALIGN.CENTER)

txt(sl, 'What TINT delivers:', 0.3, 4.65, 12, 0.4, size=14, bold=True, color=BODY)
tint_pts = [
    '— Real-time AR lipstick, eyeshadow and foundation try-on within creator posts',
    '— AI skin-tone analysis with personalised product recommendations',
    '— Embeds directly in influencer content and brand pages — no app download required',
]
y = 5.15
for pt in tint_pts:
    txt(sl, pt, 0.3, y, 12.5, 0.42, size=13, color=SECONDARY)
    y += 0.48

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0, 0.07, 7.5, LILAC)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '05 — MARKET OPPORTUNITY', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'A market ready to be shaped.', 0.3, 0.7, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

mkt = [
    ('$2.3B',  'Global AI Beauty\nPersonalization\nMarket 2026'),
    ('$2.03B', 'Influencer Platform\nMarket size\nby 2031'),
    ('45.3%',  'Instagram\'s share of\nHK influencer\nmarketing 2025'),
    ('68%',    'Beauty purchases\ninfluenced by\nsocial discovery'),
]
for i, (val, label) in enumerate(mkt):
    x = 0.3 + i * 3.25
    rect(sl, x, 1.85, 3.0, 2.8, LAV_WASH)
    rect(sl, x, 1.85, 3.0, 0.07, LILAC)
    txt(sl, val, x + 0.1, 2.05, 2.8, 0.85,
        size=30, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font='Montserrat')
    txt(sl, label, x + 0.1, 2.95, 2.8, 1.4,
        size=12, color=BODY, align=PP_ALIGN.CENTER)

rect(sl, 0.3, 5.0, 12.73, 1.8, LAV_WASH)
rect(sl, 0.3, 5.0, 12.73, 0.07, LILAC)
txt(sl, 'HK Market Insight', 0.5, 5.15, 12, 0.38, size=10, bold=True, color=MUTED)
txt(sl, 'Instagram leads HK influencer marketing at 45.3% — but Xiaohongshu (XHS) is the fastest growing platform in HK beauty. Simplex-ity is positioned across both. TikTok Shop grew beauty sales 108% in 2025. 68% of global beauty purchases are influenced at the discovery stage.',
    0.5, 5.6, 12.3, 1.1, size=13, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT BRANDS GET
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, LAV_WASH)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '06 — WHAT YOU GET', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)

txt(sl, 'Simplex-ity brand partnership — included.', 0.3, 0.7, 12, 0.65,
    size=30, bold=True, color=BODY, font='Montserrat')

deliverables = [
    ('Creator Matching',       'AI-curated creators whose audience matches your target customer — selected, not scraped.'),
    ('TINT Try-On Integration','Virtual makeup try-on embedded in posts and your brand page. No extra tech needed.'),
    ('Multi-Platform Campaigns','Instagram, XHS, TikTok — unified tracking, single dashboard.'),
    ('Real-Time Dashboard',    'Live campaign data — views, conversions, commission paid, delivery status.'),
    ('Reliability Guarantee',  'Every creator pre-vetted. Reliability scores mean no ghost deliveries.'),
    ('Full Campaign Report',   'Post-campaign ROI analysis with clear next-campaign recommendations.'),
]
cols = [(0.3, 2.1), (6.8, 2.1), (0.3, 3.7), (6.8, 3.7), (0.3, 5.3), (6.8, 5.3)]
for i, (title, desc) in enumerate(deliverables):
    x, y = cols[i]
    rect(sl, x, y, 6.2, 1.35, WHITE)
    rect(sl, x, y, 0.06, 1.35, VIOLET)
    txt(sl, title, x + 0.22, y + 0.12, 5.8, 0.42, size=14, bold=True, color=VIOLET, font='Montserrat')
    txt(sl, desc, x + 0.22, y + 0.6, 5.7, 0.65, size=12, color=BODY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — MEMBERSHIP TIERS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, WHITE)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
rect(sl, 0, 0, 0.07, 7.5, LILAC)
logo(sl, 0.3, 0.18, 1.8)
txt(sl, '07 — BRAND MEMBERSHIP', 0.3, 0.22, 5, 0.3, size=9, bold=True, color=MUTED)
txt(sl, 'Three tiers. One platform.', 0.3, 0.7, 12, 0.65,
    size=34, bold=True, color=BODY, font='Montserrat')

tiers = [
    ('Essential',   'Entry-level brands\n& new launches',
     ['1 campaign / quarter', 'Up to 5 creators', 'TINT try-on (basic)', 'Performance report'],
     False, 'HK$400 / month'),
    ('Professional','Scaling brands\n& seasonal pushes',
     ['2 campaigns / quarter', 'Up to 15 creators', 'Full TINT integration', 'Real-time dashboard', 'Priority matching'],
     True, 'HK$900 / month'),
    ('Enterprise',  'Established brands\n& multi-channel',
     ['Unlimited campaigns', 'Unlimited creators', 'Custom TINT features', 'Dedicated manager', 'Monthly strategy call'],
     False, 'HK$2,500 / month'),
]
x_pos = [0.4, 4.8, 9.2]
for i, (name, target, feats, featured, price) in enumerate(tiers):
    x = x_pos[i]
    bg_col = SOFT_LILAC if featured else LAV_WASH
    rect(sl, x, 1.85, 3.9, 5.4, bg_col)
    rect(sl, x, 1.85, 3.9, 0.07, VIOLET if featured else LILAC)
    txt(sl, name, x + 0.15, 2.0, 3.6, 0.5,
        size=20, bold=True, color=VIOLET if featured else BODY,
        align=PP_ALIGN.CENTER, font='Montserrat')
    txt(sl, target, x + 0.15, 2.58, 3.6, 0.62,
        size=12, color=SECONDARY, align=PP_ALIGN.CENTER, italic=True)
    yf = 3.32
    for feat in feats:
        txt(sl, '✓  ' + feat, x + 0.25, yf, 3.4, 0.38, size=12, color=BODY)
        yf += 0.4
    rect(sl, x + 0.2, 6.75, 3.5, 0.38, VIOLET)
    txt(sl, price, x + 0.2, 6.77, 3.5, 0.34,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / CTA
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, BODY)
rect(sl, 0, 0, 13.33, 0.07, VIOLET)
# Soft lilac right panel
rect(sl, 7.8, 0, 5.53, 7.5, RGBColor(0x2a, 0x1a, 0x4a))

logo(sl, 0.5, 0.3, 2.8)

txt(sl, 'Let\'s build something', 0.5, 1.4, 7, 0.9,
    size=40, bold=True, color=WHITE, font='Montserrat')
txt(sl, 'beautiful together.', 0.5, 2.2, 7, 0.9,
    size=40, bold=True, color=LILAC, font='Montserrat')

txt(sl, 'Join Simplex-ity as a brand partner.', 0.5, 3.3, 7, 0.5,
    size=18, color=SOFT_LILAC, italic=True)

rect(sl, 0.5, 4.2, 6.5, 0.06, LILAC)

contact = [
    'kieran@5senses.global',
    'Room 1608, 16/F APEC Plaza, 49 Hoi Yuen Road, Kwun Tong, HK',
    'SIMPLEX-ITY by 5SENSESBEAUTY LIMITED',
    'simplex-ity.fluentlab.co',
]
y = 4.4
for c in contact:
    txt(sl, c, 0.5, y, 7, 0.42, size=13, color=SOFT_LILAC)
    y += 0.5

txt(sl, 'ONE STOP ASIAN BEAUTY  :  SIMPLIFY TO AMPLIFY',
    8.1, 3.3, 5, 0.5, size=11, bold=True, color=LILAC,
    align=PP_ALIGN.CENTER)
txt(sl, 'Discover Asian beauty\nwith more confidence.',
    8.0, 4.0, 5.1, 1.5, size=22, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font='Montserrat')

prs.save('/app/SIMPLEX-ITY_Brand_Pitch_v2.pptx')
print('Done!')
